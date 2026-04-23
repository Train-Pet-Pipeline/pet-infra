"""Generate Train-Pet-Pipeline-Technical-Overview.pptx next to OVERVIEW.md.

Design brief: professional technical PPT — every slide has a diagram /
table / chart; bullets are supporting, not primary. Consistent color
palette, visual grid, typography hierarchy. Mirrors the 9-chapter
architecture.md template across 9 repos.

Regenerate:
    python pet-infra/docs/architecture/generate_technical_overview_ppt.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------

# Slide size — widescreen 16:9 in EMUs (1 inch = 914400 EMU)
SLIDE_W = Emu(12192000)  # 13.333 in
SLIDE_H = Emu(6858000)   #  7.5 in

# Color palette (hex → RGB)
NAVY        = RGBColor(0x1A, 0x2A, 0x44)  # primary dark
BLUE        = RGBColor(0x36, 0x7B, 0xD6)  # accent
TEAL        = RGBColor(0x26, 0xA2, 0x9C)  # success / stable
AMBER       = RGBColor(0xED, 0xA8, 0x3E)  # warning / deferred
CORAL       = RGBColor(0xE0, 0x5E, 0x5E)  # breaking / removed
GRAY_700    = RGBColor(0x3C, 0x40, 0x4C)
GRAY_500    = RGBColor(0x78, 0x78, 0x82)
GRAY_300    = RGBColor(0xC8, 0xCC, 0xD2)
GRAY_100    = RGBColor(0xF0, 0xF1, 0xF4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

# Repo-specific accent hues (each ~10% darker variant of a teal/blue family)
REPO_COLORS = {
    "pet-schema":     RGBColor(0x43, 0x5E, 0xB8),
    "pet-infra":      RGBColor(0x2F, 0x6A, 0xB8),
    "pet-data":       RGBColor(0x29, 0x85, 0xA8),
    "pet-annotation": RGBColor(0x26, 0xA2, 0x9C),
    "pet-train":      RGBColor(0x5A, 0xA3, 0x5E),
    "pet-eval":       RGBColor(0x8E, 0xA4, 0x3A),
    "pet-quantize":   RGBColor(0xC2, 0x8E, 0x2E),
    "pet-ota":        RGBColor(0xC2, 0x5E, 0x3E),
    "pet-id":         RGBColor(0x7A, 0x57, 0xA8),
}

# Fonts
FONT_SANS = "Helvetica"
FONT_MONO = "Menlo"


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h,
             text, *,
             size=14, color=NAVY, bold=False, italic=False,
             font=FONT_SANS, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, wrap=True):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.color.rgb = color
        f.bold = bold
        f.italic = italic
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                size=14, color=NAVY, bullet_color=BLUE,
                line_spacing=1.35, font=FONT_SANS):
    """Each item is str OR (main_str, sub_str_italic)."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Bullet mark as a small colored square
        r_mark = p.add_run()
        r_mark.text = "▪  "
        r_mark.font.name = font
        r_mark.font.size = Pt(size)
        r_mark.font.color.rgb = bullet_color
        r_mark.font.bold = True
        # Body
        if isinstance(item, tuple):
            main, sub = item
            r_main = p.add_run()
            r_main.text = main
            r_main.font.name = font
            r_main.font.size = Pt(size)
            r_main.font.color.rgb = color
            r_sub = p.add_run()
            r_sub.text = "  — " + sub
            r_sub.font.name = font
            r_sub.font.size = Pt(size - 1)
            r_sub.font.color.rgb = GRAY_500
            r_sub.font.italic = True
        else:
            r = p.add_run()
            r.text = item
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, text, *,
            fill=WHITE, border=GRAY_300, text_color=NAVY,
            size=13, bold=False, font=FONT_SANS, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
            align=PP_ALIGN.CENTER, border_w=0.75, radius=None):
    s = slide.shapes.add_shape(shape, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(border_w)
    s.shadow.inherit = False
    # If rounded rect, reduce corner radius for a subtler look
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE and hasattr(s, "adjustments"):
        try:
            s.adjustments[0] = radius if radius is not None else 0.12
        except IndexError:
            pass
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Emu(40000)
    tf.margin_top = tf.margin_bottom = Emu(30000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.2
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = text_color
        r.font.bold = bold
    return s


def add_arrow_right(slide, x, y, w, h=None, color=BLUE):
    """Right-facing arrow shape for pipeline flows."""
    if h is None:
        h = 200_000  # slim
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(x), Emu(y), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_header_bar(slide, title, subtitle=None, *, accent=BLUE, section_num=None):
    """Consistent top bar: section number tab + title."""
    # Top accent strip
    add_rect(slide, 0, 0, SLIDE_W, 80_000, fill=accent)
    # Section number tab
    if section_num:
        add_text(slide, 460_000, 200_000, 2_000_000, 350_000,
                 section_num, size=11, color=accent, bold=True, font=FONT_SANS)
    # Title
    add_text(slide, 460_000, 420_000, SLIDE_W.emu - 920_000, 600_000,
             title, size=26, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, 460_000, 1_020_000, SLIDE_W.emu - 920_000, 400_000,
                 subtitle, size=13, color=GRAY_500, italic=True)
    # Divider line below header
    add_rect(slide, 460_000, 1_460_000, SLIDE_W.emu - 920_000, 15_000,
             fill=GRAY_300)


def add_footer(slide, repo=None, slide_num=None, total=None):
    # Bottom-footer bar. If repo passed, reserve leftmost ~2" for chip+name;
    # otherwise main footer text starts at the normal left margin.
    y = SLIDE_H.emu - 440_000
    main_x = 460_000
    if repo:
        color = REPO_COLORS.get(repo, BLUE)
        # 5 EMU-wide accent stripe flush-left
        add_rect(slide, 460_000, y - 30_000, 50_000, 260_000, fill=color)
        # Repo name immediately right of stripe
        add_text(slide, 570_000, y - 20_000, 1_800_000, 260_000,
                 repo, size=10, color=color, bold=True, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Main footer text starts past the chip area
        main_x = 2_500_000
    add_text(slide, main_x, y, SLIDE_W.emu - main_x - 500_000, 300_000,
             "Train-Pet-Pipeline · Technical Overview · matrix 2026.10-ecosystem-cleanup",
             size=9, color=GRAY_500, anchor=MSO_ANCHOR.MIDDLE)


# ----------------------------------------------------------------------------
# Slide builders
# ----------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]  # blank


def new_slide():
    return prs.slides.add_slide(BLANK)


# ---- Slide 1: Title ----

def slide_title():
    s = new_slide()
    # Full-bleed navy panel
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    # Accent strip
    add_rect(s, 0, 4_700_000, SLIDE_W, 8_000, fill=BLUE)
    # Eyebrow
    add_text(s, 820_000, 2_000_000, 8_000_000, 400_000,
             "TECHNICAL OVERVIEW · TRAIN-PET-PIPELINE",
             size=12, color=BLUE, bold=True, font=FONT_SANS)
    # Title
    add_text(s, 820_000, 2_460_000, 10_000_000, 1_100_000,
             "智能宠物喂食器 AI 流水线",
             size=44, color=WHITE, bold=True)
    # Sub-title
    add_text(s, 820_000, 3_550_000, 10_000_000, 600_000,
             "9-repo monorepo · VLM + Audio + On-device Deployment",
             size=18, color=GRAY_300, italic=True)
    # Meta strip
    add_text(s, 820_000, 5_050_000, 10_000_000, 400_000,
             "matrix 2026.10-ecosystem-cleanup",
             size=13, color=BLUE, bold=True, font=FONT_MONO)
    add_text(s, 820_000, 5_420_000, 10_000_000, 400_000,
             "Post-Phase-10 ecosystem optimization closeout · 2026-04-23",
             size=12, color=GRAY_500)
    # Decorative corner mark — thicker accent lines to be visible on projection
    add_rect(s, SLIDE_W.emu - 2_100_000, 800_000, 1_400_000, 30_000, fill=BLUE)
    add_rect(s, SLIDE_W.emu - 2_100_000, 900_000, 700_000, 15_000, fill=GRAY_500)


# ---- Slide 2: 项目概览 ----

# ---- Slide: Pipeline dataflow diagram (hero) ----

def slide_pipeline_flow():
    s = new_slide()
    add_header_bar(s, "数据流水线",
                   "契约 → 数据 → 标注 → 训练 → 评估 → 量化 → 发布",
                   section_num="§ 0.1 · Pipeline Flow")

    chain = [
        ("pet-schema", "契约"),
        ("pet-data", "采集\n清洗"),
        ("pet-annotation", "4 范式\n打标"),
        ("pet-train", "SFT+DPO\n+Audio"),
        ("pet-eval", "8 指标\n+gate"),
        ("pet-quantize", "RKLLM\n+RKNN"),
        ("pet-ota", "canary\nrollout"),
    ]
    # Horizontal row of labeled boxes
    box_w = 1_340_000
    box_h = 1_050_000
    gap = 140_000
    row_y = 2_300_000
    start_x = (SLIDE_W.emu - (box_w * len(chain) + gap * (len(chain) - 1))) // 2
    for i, (name, role) in enumerate(chain):
        x = start_x + i * (box_w + gap)
        color = REPO_COLORS[name]
        # Box
        add_rect(s, x, row_y, box_w, box_h, fill=WHITE, line=color, line_w=1.5)
        # Top colored strip
        add_rect(s, x, row_y, box_w, 80_000, fill=color)
        # Name (mono)
        add_text(s, x, row_y + 170_000, box_w, 340_000,
                 name, size=12, color=color, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        # Role
        add_text(s, x, row_y + 540_000, box_w, 480_000,
                 role, size=13, color=NAVY, align=PP_ALIGN.CENTER, line_spacing=1.2)
        # Arrow
        if i < len(chain) - 1:
            arrow_x = x + box_w
            arrow_y = row_y + box_h // 2 - 100_000
            add_arrow_right(s, arrow_x + 10_000, arrow_y, gap - 20_000, 200_000, color=GRAY_500)

    # Side-branch: pet-infra runtime — placed below header divider (y≥1.55M)
    infra_color = REPO_COLORS["pet-infra"]
    infra_x = start_x + (box_w * 3 + gap * 3)
    infra_y = 1_620_000
    infra_h = 580_000
    add_box(s, infra_x - box_w // 2, infra_y, box_w * 2 + gap, infra_h,
            "pet-infra · 共享运行时 · 7 registries · orchestrator",
            fill=GRAY_100, border=infra_color, text_color=NAVY, size=12, bold=True)
    # Connector down to the main chain (thin infra-colored line)
    conn_x = infra_x + box_w // 2 + gap // 2 + 10_000
    add_rect(s, conn_x, infra_y + infra_h, 8_000, row_y - infra_y - infra_h, fill=infra_color)

    # Side-branch: pet-id independent
    pid_color = REPO_COLORS["pet-id"]
    pid_y = row_y + box_h + 500_000
    add_box(s, 660_000, pid_y, SLIDE_W.emu - 1_320_000, 550_000,
            "pet-id · 独立 CLI 工具 · 零 pet-* 运行时依赖 · PetCard registry + petid CLI",
            fill=WHITE, border=pid_color, text_color=pid_color, size=12, bold=True)

    # Legend
    legend_y = pid_y + 750_000
    add_text(s, 660_000, legend_y, 4_000_000, 280_000,
             "实线方块 = 流水线节点 · 虚框 = 独立工具 · pet-infra 作 runtime 贯穿",
             size=10, color=GRAY_500, italic=True)

    add_footer(s, slide_num=3)


# ---- Slide 4: Version matrix ----

def slide_version_matrix():
    s = new_slide()
    add_header_bar(s, "最终版本表",
                   "compatibility_matrix.yaml · 2026.10-ecosystem-cleanup · 9 仓全部 tag 已推到 main",
                   section_num="§ 0.2 · Versions")

    rows = [
        ("pet-schema",     "v3.2.1", "链首 · 无上游",                "Phase 5 additive: SFTSample + DPOSample"),
        ("pet-infra",      "v2.6.0", "β peer-dep chain",             "Phase 2 compose merge + StageRunner DRY"),
        ("pet-data",       "v1.3.0", "β peer-dep + α for schema",    "Phase 3 ingester_name / default_provenance split"),
        ("pet-annotation", "v2.1.1", "β peer-dep",                    "Phase 4+5 α exporter + F11 validator"),
        ("pet-train",      "v2.0.2", "β peer-dep",                    "Phase 5 β dual guard + JSONL consumer"),
        ("pet-eval",       "v2.3.0", "β + runtime peer for quantize", "Phase 6 ~533 LOC dead-code removal"),
        ("pet-quantize",   "v2.1.0", "β peer-dep (migrated)",         "Phase 7 7A pet-infra hardpin → peer-dep"),
        ("pet-ota",        "v2.2.0", "β peer-dep · signing=optional", "Phase 8 8A + 8B + no-hardcode"),
        ("pet-id",         "v0.2.0", "独立 · 无 pet-* deps",          "Phase 9 first CI + spec §5.2"),
    ]

    # Table-like styled rows
    col_x = [660_000, 2_860_000, 4_200_000, 6_500_000]
    col_w = [2_150_000, 1_300_000, 2_250_000, 4_400_000]
    header_y = 1_800_000
    row_h = 380_000

    # Header
    for i, head in enumerate(["仓库", "版本", "依赖形态", "本轮关键交付"]):
        add_rect(s, col_x[i], header_y, col_w[i], row_h, fill=NAVY)
        add_text(s, col_x[i] + 100_000, header_y, col_w[i] - 100_000, row_h,
                 head, size=12, color=WHITE, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Rows
    for ri, (name, ver, mode, delta) in enumerate(rows):
        y = header_y + row_h + ri * row_h
        fill = GRAY_100 if ri % 2 else WHITE
        for ci, w in enumerate(col_w):
            add_rect(s, col_x[ci], y, w, row_h, fill=fill, line=GRAY_300, line_w=0.4)
        # Repo color dot
        color = REPO_COLORS[name]
        add_rect(s, col_x[0] + 80_000, y + row_h // 2 - 80_000, 160_000, 160_000,
                 fill=color, shape=MSO_SHAPE.OVAL)
        add_text(s, col_x[0] + 300_000, y, col_w[0] - 320_000, row_h,
                 name, size=11, color=NAVY, bold=True, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col_x[1] + 100_000, y, col_w[1] - 100_000, row_h,
                 ver, size=11, color=BLUE, bold=True, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col_x[2] + 100_000, y, col_w[2] - 100_000, row_h,
                 mode, size=10.5, color=GRAY_700,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col_x[3] + 100_000, y, col_w[3] - 100_000, row_h,
                 delta, size=10, color=GRAY_500, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, slide_num=4)


# ---- Section divider helper ----

def slide_section_divider(num, title, subtitle, page_range=None):
    """page_range kept for back-compat but no longer rendered; slide numbers
    have drifted since the pitch front-matter was added. Section tag is
    sufficient wayfinding."""
    s = new_slide()
    # Dark hero panel
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    # Accent vertical bar
    add_rect(s, 820_000, 2_000_000, 14_000, 2_800_000, fill=BLUE)
    # Section number
    add_text(s, 1_000_000, 2_000_000, 8_000_000, 550_000,
             f"SECTION {num}", size=13, color=BLUE, bold=True, font=FONT_SANS)
    # Title
    add_text(s, 1_000_000, 2_500_000, 10_000_000, 1_100_000,
             title, size=40, color=WHITE, bold=True)
    # Subtitle
    add_text(s, 1_000_000, 3_750_000, 10_000_000, 500_000,
             subtitle, size=16, color=GRAY_300, italic=True)


# ---- Repo "position" diagram helper (highlight current node in the 9-chain) ----

def draw_position_strip(s, y, current):
    """Mini 9-box strip at given y, highlighting current repo.

    pet-annotation is shown as `pet-annot` in this compact strip so the
    longest name doesn't line-wrap and break row rhythm. Display only —
    the underlying repo identity remains `pet-annotation`.
    """
    chain = ["pet-schema", "pet-data", "pet-annotation", "pet-train",
             "pet-eval", "pet-quantize", "pet-ota"]
    display = {"pet-annotation": "pet-annot"}
    box_w = 950_000
    box_h = 360_000
    gap = 60_000
    total_w = len(chain) * box_w + (len(chain) - 1) * gap
    start_x = (SLIDE_W.emu - total_w) // 2
    for i, name in enumerate(chain):
        x = start_x + i * (box_w + gap)
        is_current = (name == current)
        color = REPO_COLORS[name] if is_current else GRAY_300
        fill = color if is_current else WHITE
        border = color if is_current else GRAY_300
        tcol = WHITE if is_current else GRAY_500
        add_rect(s, x, y, box_w, box_h,
                 fill=fill, line=border, line_w=1.2 if is_current else 0.6)
        add_text(s, x, y, box_w, box_h, display.get(name, name),
                 size=10, color=tcol, bold=is_current,
                 font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(chain) - 1:
            add_rect(s, x + box_w + 5_000, y + box_h // 2 - 5_000,
                     gap - 10_000, 10_000, fill=GRAY_300)
    # pet-infra band
    band_y = y + box_h + 120_000
    add_rect(s, start_x, band_y, total_w, 140_000,
             fill=REPO_COLORS["pet-infra"] if current == "pet-infra" else GRAY_100,
             line=REPO_COLORS["pet-infra"],
             line_w=1.2 if current == "pet-infra" else 0.6)
    add_text(s, start_x, band_y, total_w, 140_000,
             "pet-infra · shared runtime",
             size=9,
             color=WHITE if current == "pet-infra" else GRAY_500,
             bold=(current == "pet-infra"),
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # pet-id chip
    chip_x = start_x + total_w - 1_400_000
    chip_y = band_y + 220_000
    is_pid = current == "pet-id"
    add_rect(s, chip_x, chip_y, 1_400_000, 180_000,
             fill=REPO_COLORS["pet-id"] if is_pid else WHITE,
             line=REPO_COLORS["pet-id"],
             line_w=1.2 if is_pid else 0.6,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, chip_x, chip_y, 1_400_000, 180_000,
             "pet-id · independent",
             size=9,
             color=WHITE if is_pid else REPO_COLORS["pet-id"],
             bold=is_pid,
             font=FONT_MONO, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---- Visual primitives for per-repo technical slides ----

def draw_tech_stack_grid(s, top_y, techs, *, card_h=900_000, cols=None, gap=220_000):
    """Render a grid of tech-stack cards.

    techs: list of (name_mono, role_desc, color). Card layout:
      [color band · name (mono, bold) · role (italic small)]
    """
    n = len(techs)
    if cols is None:
        cols = 2 if n <= 4 else 3
    rows = (n + cols - 1) // cols
    usable = SLIDE_W.emu - 1_320_000
    card_w = (usable - gap * (cols - 1)) // cols
    start_x = 660_000
    for i, (name, role, color) in enumerate(techs):
        r = i // cols
        c = i % cols
        x = start_x + c * (card_w + gap)
        y = top_y + r * (card_h + gap)
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_300, line_w=0.75)
        # Left accent band
        add_rect(s, x, y, 80_000, card_h, fill=color)
        # Name (mono bold)
        add_text(s, x + 220_000, y + 160_000, card_w - 320_000, 440_000,
                 name, size=14, color=NAVY, bold=True, font=FONT_MONO)
        # Role description
        add_text(s, x + 220_000, y + 560_000, card_w - 320_000, card_h - 620_000,
                 role, size=11, color=GRAY_700, italic=True, line_spacing=1.4)


def draw_pipeline_flow(s, y, stages, *, arrow_color=GRAY_500, box_h=1_100_000):
    """Render a horizontal box+arrow pipeline.

    stages: list of (label_main, label_sub, color). Drawn with consistent
    box width (auto-computed) and a right-arrow between each pair.
    """
    n = len(stages)
    usable = SLIDE_W.emu - 1_320_000
    gap = 200_000
    box_w = (usable - gap * (n - 1)) // n
    start_x = 660_000
    for i, (label, sub, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, y, box_w, box_h, fill=WHITE, line=color, line_w=1.2)
        add_rect(s, x, y, box_w, 70_000, fill=color)
        add_text(s, x + 60_000, y + 200_000, box_w - 120_000, 440_000,
                 label, size=13, color=NAVY, bold=True,
                 font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, x + 120_000, y + 680_000, box_w - 240_000, box_h - 720_000,
                 sub, size=10, color=GRAY_500, align=PP_ALIGN.CENTER,
                 line_spacing=1.4)
        if i < n - 1:
            add_arrow_right(s, x + box_w + 20_000,
                            y + box_h // 2 - 100_000,
                            gap - 40_000, 200_000, color=arrow_color)


def draw_tech_header(s, repo, title, subtitle):
    color = REPO_COLORS[repo]
    add_header_bar(s, title, subtitle, section_num=f"§ {repo}", accent=color)


# ============================================================================
# Per-repo technical slides — tech stack + flow + design takeaway
# ============================================================================

# ---- pet-schema (2 slides) ----

def slide_schema_tech():
    s = new_slide()
    draw_tech_header(s, "pet-schema",
                     "pet-schema · 契约根节点",
                     "链首零上游 · 8 仓通过 `import pet_schema` 消费")
    # 2-col: left = tech stack, right = 主要数据模型一览
    top_y = 1_750_000
    techs = [
        ("Pydantic v2",      "extra='forbid' 严格校验\n所有跨仓数据模型",  BLUE),
        ("Alembic",          "DB schema 演化\n历史文件 immutable，只加",  TEAL),
        ("JSON Schema",      "运行时 VLM 输出校验\nPetFeederEvent v1.0",  AMBER),
        ("HF datasets",      "Features 适配器\nSFT / DPO JSONL 兼容",     CORAL),
    ]
    # Left half grid
    usable = (SLIDE_W.emu - 1_320_000) // 2 - 120_000
    card_w = (usable - 200_000) // 2
    for i, (name, role, color) in enumerate(techs):
        r, c = i // 2, i % 2
        x = 660_000 + c * (card_w + 200_000)
        y = top_y + r * (1_100_000 + 200_000)
        add_rect(s, x, y, card_w, 1_100_000, fill=WHITE, line=GRAY_300, line_w=0.75)
        add_rect(s, x, y, 80_000, 1_100_000, fill=color)
        add_text(s, x + 220_000, y + 160_000, card_w - 320_000, 420_000,
                 name, size=13, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, x + 220_000, y + 560_000, card_w - 320_000, 480_000,
                 role, size=10.5, color=GRAY_700, italic=True, line_spacing=1.45)

    # Right half: 主要 data model 按类别分组
    right_x = 660_000 + usable + 280_000
    right_w = SLIDE_W.emu - right_x - 660_000
    add_rect(s, right_x, top_y, right_w, 90_000, fill=REPO_COLORS["pet-schema"])
    add_text(s, right_x + 160_000, top_y + 160_000, right_w, 440_000,
             "核心数据模型", size=14, color=REPO_COLORS["pet-schema"], bold=True)
    groups = [
        ("训练契约",     ["Sample", "VisionSample / AudioSample",
                         "SFTSample / ShareGPTSFTSample", "DPOSample"]),
        ("编排 · 产物", ["ExperimentRecipe", "ModelCard",
                         "EdgeArtifact + QuantConfig", "DeploymentStatus"]),
        ("运行时",      ["validator.validate_output()",
                         "render_prompt(schema_version)",
                         "adapters/hf_features.py"]),
    ]
    # Flow layout per group (avoids category/items row-height clipping)
    cursor_y = top_y + 640_000
    cat_gap = 180_000
    item_h = 260_000
    for cat, items in groups:
        add_text(s, right_x + 160_000, cursor_y, right_w - 200_000, 260_000,
                 cat, size=11, color=GRAY_500, bold=True)
        cursor_y += 320_000
        for item in items:
            add_text(s, right_x + 220_000, cursor_y, right_w - 260_000, item_h,
                     item, size=10, color=NAVY, font=FONT_MONO)
            cursor_y += item_h
        cursor_y += cat_gap
    add_footer(s, repo="pet-schema")


def slide_schema_flow():
    s = new_slide()
    draw_tech_header(s, "pet-schema",
                     "pet-schema · 契约传播 + Alembic 演化",
                     "SCHEMA_VERSION 常量 + Alembic head 一致扩散到 8 仓下游")
    # Hub: pet-schema box centered
    hub_w = 3_400_000
    hub_h = 1_400_000
    hub_x = (SLIDE_W.emu - hub_w) // 2
    hub_y = 1_900_000
    color = REPO_COLORS["pet-schema"]
    add_rect(s, hub_x, hub_y, hub_w, hub_h, fill=color)
    add_text(s, hub_x, hub_y + 200_000, hub_w, 450_000,
             "pet-schema", size=22, color=WHITE, bold=True, font=FONT_MONO,
             align=PP_ALIGN.CENTER)
    add_text(s, hub_x, hub_y + 700_000, hub_w, 400_000,
             "SCHEMA_VERSION = 3.2.1", size=13, color=WHITE, bold=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, hub_x, hub_y + 1_020_000, hub_w, 350_000,
             "extra='forbid' · Alembic head: 013_schema_v3_2_1",
             size=10, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    # 8 downstream: 4 left + 4 right
    consumers = [
        ("pet-data",        "VisionSample · AudioSample\nSourceType"),
        ("pet-annotation",  "4 paradigm tables\nSFTSample · DPOSample"),
        ("pet-train",       "SFT/DPO JSONL 契约\nModelCard output"),
        ("pet-eval",        "ModelCard.metrics\nvalidate_output"),
        ("pet-quantize",    "EdgeArtifact\nQuantConfig · calib_dir"),
        ("pet-ota",         "DeploymentStatus\nModelCard.deployment_history"),
        ("pet-infra",       "Recipe · StageRunner 基础契约"),
        ("render_prompt",   "train / infer 同源 prompt"),
    ]
    left_consumers = consumers[:4]
    right_consumers = consumers[4:]
    card_w = 2_700_000
    card_h = 680_000
    gap = 100_000
    # left column
    for i, (name, what) in enumerate(left_consumers):
        x = 660_000
        y = 1_850_000 + i * (card_h + gap)
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_300, line_w=0.6)
        add_rect(s, x, y, 60_000, card_h, fill=REPO_COLORS.get(name, GRAY_500))
        add_text(s, x + 160_000, y + 60_000, card_w - 220_000, 280_000,
                 name, size=11, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, x + 160_000, y + 340_000, card_w - 220_000, card_h - 380_000,
                 what, size=9.5, color=GRAY_500, line_spacing=1.4)
        # Arrow from hub to this card
        add_arrow_right(s, x + card_w, y + card_h // 2 - 50_000,
                        hub_x - (x + card_w) - 40_000, 100_000, color=color)
    # right column (arrows point FROM hub to card → use LEFT_ARROW? python-pptx lacks directly — use a thin bar + label visual)
    for i, (name, what) in enumerate(right_consumers):
        x = SLIDE_W.emu - 660_000 - card_w
        y = 1_850_000 + i * (card_h + gap)
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_300, line_w=0.6)
        add_rect(s, x + card_w - 60_000, y, 60_000, card_h, fill=REPO_COLORS.get(name, GRAY_500))
        add_text(s, x + 160_000, y + 60_000, card_w - 220_000, 280_000,
                 name, size=11, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, x + 160_000, y + 340_000, card_w - 220_000, card_h - 380_000,
                 what, size=9.5, color=GRAY_500, line_spacing=1.4)
        # Connector
        add_arrow_right(s, hub_x + hub_w + 20_000, y + card_h // 2 - 50_000,
                        x - (hub_x + hub_w) - 40_000, 100_000, color=color)

    add_footer(s, repo="pet-schema")


# ---- pet-infra (3 slides) ----

def slide_infra_tech():
    s = new_slide()
    draw_tech_header(s, "pet-infra",
                     "pet-infra · 核心技术栈",
                     "共享运行时 — 其他 8 仓都以 pet-infra 为 peer-dep")
    techs = [
        ("mmengine Registry",  "7 registries 基座\nTRAINERS / EVALUATORS / CONVERTERS / METRICS / DATASETS / STORAGE / OTA", BLUE),
        ("Hydra + hydra-zen",  "defaults-list 组合 recipe\ncomposable + override-friendly", TEAL),
        ("entry_points",       "plugin 发现机制\ngroup=\"pet_infra.plugins\"",                AMBER),
        ("networkx",           "Stage DAG + topological_sort\nresume-from-cache 确定性",      REPO_COLORS["pet-infra"]),
        ("ClearML",            "唯一实验追踪\nW&B 已 Phase 4 物理移除",                        CORAL),
        ("click",              "CLI 子命令\npet run / replay / sweep",                         GRAY_700),
    ]
    draw_tech_stack_grid(s, 1_800_000, techs, cols=3, card_h=1_600_000, gap=220_000)
    add_footer(s, repo="pet-infra")


def slide_infra_flow():
    s = new_slide()
    draw_tech_header(s, "pet-infra",
                     "pet-infra · Recipe → ModelCard 流水线",
                     "compose_recipe() → DAG → StageRunner loop → ModelCard + ClearML")
    stages = [
        ("recipe.yaml",      "Hydra defaults-list\n+ overrides",                     GRAY_500),
        ("compose_recipe()", "hydra-zen 解析\n+ variations 展开",                      BLUE),
        ("build_dag()",      "networkx topological_sort\n+ stage.depends_on",        TEAL),
        ("StageRunner",      "_load_stage_kwargs\n+ registry.build(**cfg)",          AMBER),
        ("plugin.run()",     "consumer repo 的\n@register_module 实现",                REPO_COLORS["pet-infra"]),
        ("ModelCard",        "pet-schema 契约\n+ ClearMLLogger push",                CORAL),
    ]
    draw_pipeline_flow(s, 2_000_000, stages, box_h=1_400_000)
    # Below: side-channel — StageCache + resume
    note_y = 3_600_000
    add_rect(s, 660_000, note_y, SLIDE_W.emu - 1_320_000, 1_300_000, fill=GRAY_100)
    add_text(s, 820_000, note_y + 160_000, 4_000_000, 400_000,
             "SIDE CHANNELS", size=11, color=BLUE, bold=True)
    sides = [
        ("StageCache (~/.pet-cache)", "card_id = hash(recipe, stage, config_sha)\n缓存命中 → 跳过执行；resume 可恢复任意 stage"),
        ("ClearMLLogger",             "每 stage start/end 自动 push\nmetrics + ModelCard 对齐可 replay"),
        ("Replay (pet run --replay)", "从历史 ModelCard 反演 config\n确定性重跑训练 run"),
    ]
    sw = (SLIDE_W.emu - 1_640_000) // 3
    for i, (title, body) in enumerate(sides):
        sx = 820_000 + i * sw
        add_text(s, sx, note_y + 560_000, sw - 120_000, 280_000,
                 title, size=11, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, sx, note_y + 840_000, sw - 120_000, 400_000,
                 body, size=10, color=GRAY_700, line_spacing=1.4)
    add_footer(s, repo="pet-infra")


def slide_infra_registries():
    s = new_slide()
    draw_tech_header(s, "pet-infra",
                     "pet-infra · 7 Registries · 插件发现",
                     "每仓通过 entry_points 注册；orchestrator 按 registry.get(type) 构造插件")
    # 7 registries — 7 equal-size cells across 2 rows (4+3)
    reg_cells = [
        ("TRAINERS",   "pet-train",              "llamafactory_sft · llamafactory_dpo · tiny_test"),
        ("EVALUATORS", "pet-eval",               "vlm · audio · quantized_vlm + 3 fusion"),
        ("METRICS",    "pet-eval",               "schema / anomaly / mood / narrative / latency / audio / kl / calibration"),
        ("CONVERTERS", "pet-quantize",           "noop · vlm_rkllm_w4a16 · audio_rknn_fp16 · vision_rknn_fp16"),
        ("DATASETS",   "pet-data · pet-quantize", "vlm / vision / audio calibration_subset"),
        ("STORAGE",    "pet-infra",              "local · s3 · http · file  (STORAGE.build(uri))"),
        ("OTA",        "pet-ota",                "local_backend · s3_backend · http_backend"),
    ]
    start_y = 1_800_000
    card_w = 2_720_000
    card_h = 1_500_000
    gap = 200_000
    row1 = reg_cells[:4]
    row2 = reg_cells[4:]
    # Row 1 (4 cells, centered)
    r1_total = 4 * card_w + 3 * gap
    r1_x0 = (SLIDE_W.emu - r1_total) // 2
    for i, (reg, owner, content) in enumerate(row1):
        x = r1_x0 + i * (card_w + gap)
        y = start_y
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=REPO_COLORS[owner.split(" ")[0]], line_w=1.2)
        add_rect(s, x, y, card_w, 70_000, fill=REPO_COLORS[owner.split(" ")[0]])
        add_text(s, x, y + 160_000, card_w, 400_000, reg, size=14, color=NAVY,
                 bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, x, y + 600_000, card_w, 280_000, f"owner · {owner}", size=9,
                 color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 160_000, y + 900_000, card_w - 320_000, card_h - 950_000,
                 content, size=9.5, color=GRAY_700, line_spacing=1.45, align=PP_ALIGN.CENTER)
    # Row 2 (3 cells, centered)
    r2_total = 3 * card_w + 2 * gap
    r2_x0 = (SLIDE_W.emu - r2_total) // 2
    for i, (reg, owner, content) in enumerate(row2):
        x = r2_x0 + i * (card_w + gap)
        y = start_y + card_h + 260_000
        owner_key = owner.split(" ")[0]
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=REPO_COLORS[owner_key], line_w=1.2)
        add_rect(s, x, y, card_w, 70_000, fill=REPO_COLORS[owner_key])
        add_text(s, x, y + 160_000, card_w, 400_000, reg, size=14, color=NAVY,
                 bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, x, y + 600_000, card_w, 280_000, f"owner · {owner}", size=9,
                 color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 160_000, y + 900_000, card_w - 320_000, card_h - 950_000,
                 content, size=9.5, color=GRAY_700, line_spacing=1.45, align=PP_ALIGN.CENTER)
    add_footer(s, repo="pet-infra")


# ---- pet-data (2 slides) ----

def slide_data_tech():
    s = new_slide()
    draw_tech_header(s, "pet-data",
                     "pet-data · 数据采集 + 清洗技术栈",
                     "7 Ingester · FFmpeg 解帧 · dHash 去重 · blur/brightness 质量闸 · SQLite + Alembic")
    techs = [
        ("FFmpeg",        "视频解帧 · 提取关键帧\n1080p @ 可配置 FPS",           BLUE),
        ("dHash (pHash)", "感知哈希去重\n相似度阈值 + 汉明距离",                 TEAL),
        ("PIL · OpenCV",  "blur 检测 + brightness 直方图\nQualityFilter 门槛",   AMBER),
        ("7 Ingester",    "YouTubeIngester · CommunityIngester · Selfshot\nMock · Dryrun · ...",  REPO_COLORS["pet-data"]),
        ("SQLite + Alembic", "frames + audio_samples 两表\n004 migration · 历史 immutable", CORAL),
        ("DATASETS plugin",  "VisionSample / AudioSample export\n通过 pet-infra registry 被下游消费", GRAY_700),
    ]
    draw_tech_stack_grid(s, 1_800_000, techs, cols=3, card_h=1_500_000, gap=220_000)
    add_footer(s, repo="pet-data")


def slide_data_flow():
    s = new_slide()
    draw_tech_header(s, "pet-data",
                     "pet-data · Ingest Pipeline",
                     "从原始 URL 到 annotation-ready Sample · 每 stage 可独立扩展")
    stages = [
        ("Ingester",      "YouTube / Community / ...\nURL → 本地 mp4",                BLUE),
        ("Frame Extract", "FFmpeg 解帧\n提取至 frames/*.jpg",                          TEAL),
        ("Dedup",         "dHash 感知哈希\n相似度 ≤ 阈值则去重",                         AMBER),
        ("Quality Filter", "blur 方差 ≥ 下界\nbrightness 分布合格",                    REPO_COLORS["pet-data"]),
        ("Anomaly Score", "弱监督打分\n标注任务优先级",                                 CORAL),
        ("SQLite Write",  "frames 表 + Alembic\nannotation_status='pending'",          GRAY_700),
    ]
    draw_pipeline_flow(s, 2_100_000, stages, box_h=1_500_000)
    # Bottom note
    note_y = 3_800_000
    add_rect(s, 660_000, note_y, SLIDE_W.emu - 1_320_000, 1_000_000, fill=GRAY_100)
    add_text(s, 820_000, note_y + 180_000, SLIDE_W.emu - 1_640_000, 300_000,
             "下游消费", size=11, color=BLUE, bold=True)
    add_text(s, 820_000, note_y + 490_000, SLIDE_W.emu - 1_640_000, 460_000,
             "pet-annotation 只读 frames 表（跨仓只读 URI 模式）· "
             "pet-train DATASETS plugin 读 VisionSample / AudioSample · "
             "dedup 铁律：禁 skip 旗标（feedback_no_manual_workaround）",
             size=11, color=GRAY_700, line_spacing=1.5)
    add_footer(s, repo="pet-data")


# ---- pet-annotation (2 slides) ----

def slide_annotation_tech():
    s = new_slide()
    draw_tech_header(s, "pet-annotation",
                     "pet-annotation · 4 范式打标架构",
                     "LLM · classifier · rule · human 并行独立存 —— 不跨 annotator reconcile")
    # 4 lanes visualized horizontally
    start_y = 1_850_000
    lane_w = 2_750_000
    lane_h = 2_700_000
    gap = 200_000
    total_w = 4 * lane_w + 3 * gap
    start_x = (SLIDE_W.emu - total_w) // 2
    lanes = [
        {
            "name": "LLM",
            "impl": "OpenAI-compatible API",
            "desc": "Qwen2-VL / GPT-4o\n结构化 JSON prompt\nrender_prompt 同源",
            "table": "LLMAnnotation",
            "color": BLUE,
        },
        {
            "name": "classifier",
            "impl": "PyTorch fine-tuned",
            "desc": "小模型监督分类\nhot-swap 可替换\nbatch inference",
            "table": "ClassifierAnnotation",
            "color": TEAL,
        },
        {
            "name": "rule",
            "impl": "programmatic",
            "desc": "确定性规则\nheuristic 基线\nboxed metrics",
            "table": "RuleAnnotation",
            "color": AMBER,
        },
        {
            "name": "human",
            "impl": "Label Studio 1.23",
            "desc": "session auth\nimport / export\nDPO pair 构造",
            "table": "HumanAnnotation",
            "color": CORAL,
        },
    ]
    for i, L in enumerate(lanes):
        x = start_x + i * (lane_w + gap)
        add_rect(s, x, start_y, lane_w, lane_h, fill=WHITE, line=L["color"], line_w=1.2)
        add_rect(s, x, start_y, lane_w, 80_000, fill=L["color"])
        add_text(s, x, start_y + 160_000, lane_w, 440_000,
                 L["name"], size=18, color=L["color"], bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, start_y + 640_000, lane_w, 300_000,
                 L["impl"], size=10.5, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
        add_text(s, x + 160_000, start_y + 1_100_000, lane_w - 320_000, 1_000_000,
                 L["desc"], size=11, color=NAVY, line_spacing=1.5, align=PP_ALIGN.CENTER)
        # Storage table pill
        add_rect(s, x + 300_000, start_y + lane_h - 500_000, lane_w - 600_000, 380_000,
                 fill=GRAY_100, line=GRAY_300, line_w=0.5)
        add_text(s, x, start_y + lane_h - 460_000, lane_w, 280_000,
                 L["table"], size=10, color=NAVY, bold=True,
                 font=FONT_MONO, align=PP_ALIGN.CENTER)
    # Caption
    add_text(s, 660_000, 4_950_000, SLIDE_W.emu - 1_320_000, 320_000,
             "每 annotator 写各自表（D4 决策 · 不做 majority-vote reconcile）· "
             "AnnotationOrchestrator 并发调度 (target × annotator)",
             size=10.5, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, repo="pet-annotation")


def slide_annotation_flow():
    s = new_slide()
    draw_tech_header(s, "pet-annotation",
                     "pet-annotation · Orchestrator + SFT/DPO Exporter",
                     "并发 fan-out → 4 独立表 → LLaMA-Factory JSONL (α 方向，spec §5.3)")
    # Fan-out diagram — 4 columns: target | orchestrator | 4 branches | exporter
    row_y = 2_100_000
    row_h = 2_800_000

    left_x = 660_000
    target_w = 1_900_000
    # target box
    add_rect(s, left_x, row_y, target_w, row_h, fill=NAVY)
    add_text(s, left_x, row_y + 360_000, target_w, 500_000,
             "target_id", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
             font=FONT_MONO)
    add_text(s, left_x + 120_000, row_y + 1_000_000, target_w - 240_000, 800_000,
             "从 pet-data frames\n(pending targets)", size=10, color=GRAY_300,
             align=PP_ALIGN.CENTER, italic=True, line_spacing=1.5)

    # orchestrator
    orch_x = left_x + target_w + 180_000
    orch_w = 2_400_000
    add_rect(s, orch_x, row_y, orch_w, row_h, fill=REPO_COLORS["pet-annotation"])
    add_text(s, orch_x, row_y + 360_000, orch_w, 500_000,
             "AnnotationOrchestrator", size=13, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, font=FONT_MONO)
    add_text(s, orch_x + 120_000, row_y + 1_000_000, orch_w - 240_000, 1_000_000,
             "并发 (target × annotator)\n组合调度\n每 annotator 独立写表",
             size=11, color=WHITE, align=PP_ALIGN.CENTER, italic=True, line_spacing=1.5)

    # Arrow target → orchestrator
    add_arrow_right(s, left_x + target_w + 20_000, row_y + row_h // 2 - 100_000,
                    140_000, 200_000, color=GRAY_500)

    # 4 branches
    branch_x = orch_x + orch_w + 260_000
    branch_w = 1_900_000
    branch_h = 520_000
    branch_gap = 140_000
    branches = [
        ("LLM",        BLUE,   "LLMAnnotation"),
        ("classifier", TEAL,   "ClassifierAnnotation"),
        ("rule",       AMBER,  "RuleAnnotation"),
        ("human",      CORAL,  "HumanAnnotation"),
    ]
    total_branches_h = 4 * branch_h + 3 * branch_gap
    by_start = row_y + (row_h - total_branches_h) // 2
    for i, (name, c, table) in enumerate(branches):
        y = by_start + i * (branch_h + branch_gap)
        add_rect(s, branch_x, y, branch_w, branch_h, fill=WHITE, line=c, line_w=1.0)
        add_rect(s, branch_x, y, 70_000, branch_h, fill=c)
        add_text(s, branch_x + 140_000, y + 60_000, branch_w - 180_000, 220_000,
                 name, size=12, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, branch_x + 140_000, y + 280_000, branch_w - 180_000, 220_000,
                 table, size=9, color=GRAY_500, font=FONT_MONO, italic=True)
        # Connector from orchestrator right edge to each branch left
        conn_src_x = orch_x + orch_w
        conn_dst_x = branch_x
        conn_y = y + branch_h // 2 - 4_000
        add_rect(s, conn_src_x, conn_y, conn_dst_x - conn_src_x, 8_000,
                 fill=REPO_COLORS["pet-annotation"])

    # Exporter box
    exp_x = branch_x + branch_w + 200_000
    exp_w = SLIDE_W.emu - exp_x - 660_000
    # Connect each branch to exporter left edge via a short horizontal line
    for i in range(4):
        y = by_start + i * (branch_h + branch_gap) + branch_h // 2 - 4_000
        add_rect(s, branch_x + branch_w, y, exp_x - (branch_x + branch_w), 8_000,
                 fill=REPO_COLORS["pet-annotation"])

    exp_h = 2_100_000
    exp_y = row_y + 240_000
    add_rect(s, exp_x, exp_y, exp_w, exp_h, fill=WHITE,
             line=REPO_COLORS["pet-annotation"], line_w=1.2)
    add_rect(s, exp_x, exp_y, exp_w, 70_000, fill=REPO_COLORS["pet-annotation"])
    add_text(s, exp_x, exp_y + 160_000, exp_w, 400_000,
             "sft_dpo.py exporter", size=13, color=NAVY, bold=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, exp_x + 160_000, exp_y + 640_000, exp_w - 320_000, 700_000,
             "to_sft_samples()  →  SFTSample JSONL\n"
             "to_dpo_pairs()    →  DPOSample JSONL",
             size=10.5, color=GRAY_700, line_spacing=1.6, font=FONT_MONO)
    # F11 box inside
    add_rect(s, exp_x + 160_000, exp_y + 1_420_000, exp_w - 320_000, 580_000,
             fill=GRAY_100)
    add_text(s, exp_x + 280_000, exp_y + 1_480_000, exp_w - 560_000, 280_000,
             "F11 producer-side validator", size=10, color=NAVY, bold=True, font=FONT_MONO)
    add_text(s, exp_x + 280_000, exp_y + 1_760_000, exp_w - 560_000, 300_000,
             "DPOSample.model_validate() per row",
             size=9, color=GRAY_500, italic=True, font=FONT_MONO)

    # Consumer band
    pt_y = exp_y + exp_h + 180_000
    add_rect(s, exp_x, pt_y, exp_w, 460_000, fill=REPO_COLORS["pet-train"])
    add_text(s, exp_x, pt_y + 110_000, exp_w, 300_000,
             "→ pet-train consumes", size=12, color=WHITE, bold=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)

    add_footer(s, repo="pet-annotation")


# ---- pet-train (2 slides) ----

def slide_train_tech():
    s = new_slide()
    draw_tech_header(s, "pet-train",
                     "pet-train · 核心技术栈",
                     "LLaMA-Factory vendored · PEFT LoRA · DPO · PANNs · F11 consumer validator")
    techs = [
        ("LLaMA-Factory v0.9.4", "Apache-2.0 vendored (plain dir)\nrun_sft / run_dpo workflow",       BLUE),
        ("PEFT LoRA",            "低秩适配微调 (r / alpha / lr)\nHF peft + Transformers",            TEAL),
        ("DPO · pref_beta=0.1",  "sft_adapter_path 为 base\n从 pet-annotation DPO pair 训练",         AMBER),
        ("PANNs MobileNetV2",    "AudioSet 527 classes\nmel-spectrogram zero-shot",                   REPO_COLORS["pet-train"]),
        ("Lazy run_sft import",  "module-load 不拉 LLaMA-Factory\n避免 transformers 脆弱上游",        CORAL),
        ("F11 JSONL validator",  "validate_sft_jsonl / validate_dpo_jsonl\n训练前严校 (双端都验)",    GRAY_700),
    ]
    draw_tech_stack_grid(s, 1_800_000, techs, cols=3, card_h=1_550_000, gap=220_000)
    add_footer(s, repo="pet-train")


def slide_train_flow():
    s = new_slide()
    draw_tech_header(s, "pet-train",
                     "pet-train · SFT + DPO 训练流水线",
                     "共享 LoRA + F11 validator · DPO 以 SFT adapter 为 base")
    # Two parallel rows: SFT pipeline, DPO pipeline
    row_h = 1_200_000
    sft_y = 1_900_000
    dpo_y = sft_y + row_h + 300_000

    def draw_row(y, title, color, stages):
        # Row label
        add_rect(s, 660_000, y, 1_100_000, row_h, fill=color)
        add_text(s, 660_000, y, 1_100_000, row_h, title, size=18, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Stages
        start_x = 660_000 + 1_100_000 + 200_000
        usable = SLIDE_W.emu - 660_000 - start_x
        n = len(stages)
        gap = 140_000
        box_w = (usable - gap * (n - 1)) // n
        for i, (label, sub) in enumerate(stages):
            x = start_x + i * (box_w + gap)
            add_rect(s, x, y, box_w, row_h, fill=WHITE, line=color, line_w=1.0)
            add_rect(s, x, y, box_w, 60_000, fill=color)
            add_text(s, x, y + 180_000, box_w, 400_000,
                     label, size=12, color=NAVY, bold=True,
                     font=FONT_MONO, align=PP_ALIGN.CENTER)
            add_text(s, x + 100_000, y + 600_000, box_w - 200_000, row_h - 660_000,
                     sub, size=10, color=GRAY_500, align=PP_ALIGN.CENTER, line_spacing=1.4)
            if i < n - 1:
                add_arrow_right(s, x + box_w + 20_000, y + row_h // 2 - 80_000,
                                gap - 40_000, 160_000, color=GRAY_500)

    draw_row(sft_y, "SFT", BLUE, [
        ("SFT JSONL",        "ShareGPT format\nfrom pet-annotation"),
        ("F11 validator",    "validate_sft_jsonl\n早失败 file:line"),
        ("run_sft(**cfg)",   "lora_r / alpha / lr\nfrom params.yaml"),
        ("LoRA adapter",     "output_dir/adapter\nPEFT safetensors"),
        ("ModelCard",        "checkpoint_uri\n+ metrics"),
    ])
    draw_row(dpo_y, "DPO", AMBER, [
        ("DPO JSONL",        "DPOSample pairs\nchosen/rejected"),
        ("F11 validator",    "validate_dpo_jsonl\npref pair schema"),
        ("run_dpo(**cfg)",   "pref_beta=0.1\nsft_adapter_path"),
        ("DPO adapter",      "基于 SFT 继续\n对齐 preference"),
        ("ModelCard",        "checkpoint_uri\n+ metrics"),
    ])
    # Bottom: tiny_test mention
    tiny_y = dpo_y + row_h + 280_000
    add_rect(s, 660_000, tiny_y, SLIDE_W.emu - 1_320_000, 450_000, fill=GRAY_100)
    add_text(s, 820_000, tiny_y + 100_000, SLIDE_W.emu - 1_640_000, 300_000,
             "tiny_test trainer · CPU-only smoke <2min · 独立 plugin · PR-gate 不需 GPU",
             size=11, color=GRAY_700, italic=True)
    add_footer(s, repo="pet-train")


# ---- pet-eval (2 slides) ----

def slide_eval_tech():
    s = new_slide()
    draw_tech_header(s, "pet-eval",
                     "pet-eval · 8 Metrics + 6 Evaluators 架构",
                     "@METRICS / @EVALUATORS 装饰注册 · rule-based fusion only · apply_gate 决策")
    # Two-col: left = 8 metrics, right = 6 evaluators
    top_y = 1_800_000
    col_gap = 300_000
    col_w = (SLIDE_W.emu - 1_320_000 - col_gap) // 2
    left_x = 660_000
    right_x = left_x + col_w + col_gap

    # Left: METRICS (8 cards in 4×2 grid inside left col)
    add_rect(s, left_x, top_y, col_w, 90_000, fill=REPO_COLORS["pet-eval"])
    add_text(s, left_x + 160_000, top_y + 160_000, col_w, 400_000,
             "@METRICS · 8", size=16, color=REPO_COLORS["pet-eval"], bold=True)
    metrics = [
        ("schema_compliance",    "compliance_rate + distribution_sum_error"),
        ("anomaly_recall",       "TPR + FPR"),
        ("mood_correlation",     "Spearman × 3 dim"),
        ("narrative_quality",    "BERTScore F1 (Chinese)"),
        ("latency",              "P50 / P95 / P99 interp"),
        ("audio_accuracy",       "overall + vomit_recall"),
        ("kl_quantization",      "fp16 ↔ quant distribution"),
        ("calibration",          "ECE · informational only"),
    ]
    m_start = top_y + 620_000
    m_h = 440_000
    for i, (name, sub) in enumerate(metrics):
        r = i // 2
        c = i % 2
        iw = (col_w - 320_000) // 2
        x = left_x + 160_000 + c * (iw + 0)
        y = m_start + r * m_h
        add_rect(s, x, y, iw, m_h - 60_000, fill=WHITE, line=GRAY_300, line_w=0.5)
        add_text(s, x + 80_000, y + 50_000, iw - 140_000, 220_000,
                 name, size=10, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, x + 80_000, y + 250_000, iw - 140_000, 140_000,
                 sub, size=8.5, color=GRAY_500, italic=True)

    # Right: EVALUATORS (3 primary + 3 fusion)
    add_rect(s, right_x, top_y, col_w, 90_000, fill=REPO_COLORS["pet-eval"])
    add_text(s, right_x + 160_000, top_y + 160_000, col_w, 400_000,
             "@EVALUATORS · 6", size=16, color=REPO_COLORS["pet-eval"], bold=True)
    primary = [
        ("vlm_evaluator",           "Qwen2-VL + PEFT LoRA merge\n→ gold set inference → metrics"),
        ("audio_evaluator",         "→ pet_train.audio.inference\n(跨仓 runtime)"),
        ("quantized_vlm_evaluator", "→ pet_quantize.rkllm_runner\n(跨仓 runtime, lazy)"),
    ]
    fusion = [
        ("single_modal_fusion",  "pass-through scores[modality]"),
        ("and_gate_fusion",      "all ≥ threshold → min else 0"),
        ("weighted_fusion",      "normalized weighted sum"),
    ]
    p_start = top_y + 620_000
    p_h = 820_000
    add_text(s, right_x + 160_000, p_start - 30_000, col_w, 240_000,
             "PRIMARY (3)", size=9, color=GRAY_500, bold=True, italic=True)
    for i, (name, sub) in enumerate(primary):
        y = p_start + 220_000 + i * p_h
        add_rect(s, right_x + 160_000, y, col_w - 320_000, p_h - 100_000,
                 fill=WHITE, line=GRAY_300, line_w=0.5)
        add_text(s, right_x + 240_000, y + 70_000, col_w - 480_000, 260_000,
                 name, size=10.5, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, right_x + 240_000, y + 340_000, col_w - 480_000, p_h - 480_000,
                 sub, size=9, color=GRAY_500, italic=True, line_spacing=1.35)
    f_start = p_start + 220_000 + 3 * p_h + 60_000
    add_text(s, right_x + 160_000, f_start - 30_000, col_w, 240_000,
             "FUSION · rule-based (3) · no learned", size=9, color=GRAY_500, bold=True, italic=True)
    for i, (name, sub) in enumerate(fusion):
        iw = (col_w - 320_000 - 120_000) // 3
        x = right_x + 160_000 + i * (iw + 60_000)
        y = f_start + 220_000
        add_rect(s, x, y, iw, 380_000, fill=WHITE, line=GRAY_300, line_w=0.5)
        add_text(s, x + 60_000, y + 40_000, iw - 120_000, 160_000,
                 name, size=9, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, x + 60_000, y + 200_000, iw - 120_000, 160_000,
                 sub, size=8, color=GRAY_500, italic=True, line_spacing=1.3)
    add_footer(s, repo="pet-eval")


def slide_eval_flow():
    s = new_slide()
    draw_tech_header(s, "pet-eval",
                     "pet-eval · 推理 → 指标 → Gate 决策",
                     "ModelCard in · ModelCard out (gate_status: passed / failed · reason)")
    stages = [
        ("input ModelCard",   "checkpoint_uri\n(from pet-train)",              GRAY_500),
        ("run_inference",     "HF + PEFT merge\ngold_set JSONL → outputs",      BLUE),
        ("_compute_metrics",  "8 registered metrics\nwarn-and-skip on sig mismatch", TEAL),
        ("apply_gate",        "min_<m> / max_<m>\nthresholds → pass/fail",      AMBER),
        ("FusionEvaluator",   "rule-based 3 策略\n可选 (若 recipe 声明)",         REPO_COLORS["pet-eval"]),
        ("output ModelCard",  "+ metrics dict\n+ gate_status + reason",          CORAL),
    ]
    draw_pipeline_flow(s, 2_000_000, stages, box_h=1_400_000)
    # Safety note
    note_y = 3_600_000
    add_rect(s, 660_000, note_y, SLIDE_W.emu - 1_320_000, 1_300_000, fill=GRAY_100)
    add_text(s, 820_000, note_y + 160_000, 5_000_000, 400_000,
             "设计要点", size=11, color=BLUE, bold=True)
    points = [
        ("_FALLBACK_OUTPUT",   "VLM 重试失败后的安全 JSON sentinel\n防止 metric 分母错"),
        ("prompt_source",      "gold_set | pet_schema 双轨\nsft_v2 / sft_v3+ 都可评"),
        ("跨仓 lazy import",   "pet_train.audio · pet_quantize.rkllm\nmodule-load 不拉 SDK"),
    ]
    pw = (SLIDE_W.emu - 1_640_000) // 3
    for i, (title, body) in enumerate(points):
        px = 820_000 + i * pw
        add_text(s, px, note_y + 580_000, pw - 120_000, 300_000,
                 title, size=11, color=NAVY, bold=True, font=FONT_MONO)
        add_text(s, px, note_y + 860_000, pw - 120_000, 400_000,
                 body, size=10, color=GRAY_700, line_spacing=1.4)
    add_footer(s, repo="pet-eval")


# ---- pet-quantize (3 slides) ----

def slide_quantize_tech():
    s = new_slide()
    draw_tech_header(s, "pet-quantize",
                     "pet-quantize · 核心技术栈",
                     "Rockchip RKNN / RKLLM SDK + torch.onnx + ADB + 签名打包")
    techs = [
        ("rkllm-toolkit",            "VLM → RKLLM W4A16\nPANDAS calibration batch",        BLUE),
        ("rknn-toolkit2",            "ViT → RKNN FP16\nAudio CNN → RKNN FP16",             TEAL),
        ("torch.onnx",               "vision encoder 导出\n中间格式",                      AMBER),
        ("ADB bridge",               "on-device execution\ntarget='rk3576' + device_id",    REPO_COLORS["pet-quantize"]),
        ("Content-address cache",    "sha256(modality|uri|n)[:16]\norchestrator resume-from-cache", CORAL),
        ("SDK-gated cluster",        "PET_ALLOW_MISSING_SDK=1\n按 SDK 分组精细降级",        GRAY_700),
    ]
    draw_tech_stack_grid(s, 1_800_000, techs, cols=3, card_h=1_550_000, gap=220_000)
    add_footer(s, repo="pet-quantize")


def slide_quantize_flow():
    s = new_slide()
    draw_tech_header(s, "pet-quantize",
                     "pet-quantize · Convert → EdgeArtifact",
                     "ModelCard 分三簇按需转换 · 缺 SDK 时 cluster-level 精细降级")
    # Layout: 3 clusters stacked vertically in middle column, ModelCard left-middle, EdgeArtifact right-middle
    br_w = 3_200_000
    br_h = 1_050_000
    br_gap = 180_000
    top_cluster_y = 1_950_000

    in_w = 2_000_000
    in_h = 1_000_000
    in_x = 660_000
    out_w = 2_200_000

    br_x = in_x + in_w + 600_000
    out_x = br_x + br_w + 600_000

    mid_cluster_y = top_cluster_y + br_h + br_gap
    bot_cluster_y = mid_cluster_y + br_h + br_gap

    # ModelCard aligned with middle cluster center
    in_y = mid_cluster_y + (br_h - in_h) // 2
    add_rect(s, in_x, in_y, in_w, in_h, fill=NAVY)
    add_text(s, in_x, in_y + 180_000, in_w, 400_000,
             "ModelCard", size=14, color=WHITE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, in_x, in_y + 580_000, in_w, 380_000,
             "checkpoint_uri\n+ calibration_batch_uri", size=9.5, color=GRAY_300,
             italic=True, align=PP_ALIGN.CENTER, line_spacing=1.35)

    clusters = [
        ("Always-available", "noop_converter · 零 SDK · CI smoke", GRAY_500, top_cluster_y),
        ("rknn cluster",     "audio_rknn_fp16 · vision_rknn_fp16\n+ calibration DATASETS", TEAL, mid_cluster_y),
        ("rkllm cluster",    "vlm_rkllm_w4a16 · + vlm_calibration_subset", CORAL, bot_cluster_y),
    ]
    for i, (name, content, color, y) in enumerate(clusters):
        add_rect(s, br_x, y, br_w, br_h, fill=WHITE, line=color, line_w=1.2)
        add_rect(s, br_x, y, br_w, 70_000, fill=color)
        add_text(s, br_x, y + 180_000, br_w, 400_000,
                 name, size=14, color=NAVY, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, br_x + 160_000, y + 620_000, br_w - 320_000, br_h - 680_000,
                 content, size=10, color=GRAY_700, line_spacing=1.4, align=PP_ALIGN.CENTER)
        # Thin connector from ModelCard right edge to cluster left edge
        conn_src_x = in_x + in_w
        conn_y = y + br_h // 2 - 4_000
        add_rect(s, conn_src_x, conn_y, br_x - conn_src_x, 8_000, fill=color)
        # Thin connector from cluster right edge to EdgeArtifact left edge
        add_rect(s, br_x + br_w, conn_y, out_x - (br_x + br_w), 8_000, fill=color)

    # Output EdgeArtifact aligned with middle cluster center
    out_h = 1_000_000
    out_y = in_y
    add_rect(s, out_x, out_y, out_w, out_h, fill=REPO_COLORS["pet-quantize"])
    add_text(s, out_x, out_y + 180_000, out_w, 400_000,
             "EdgeArtifact", size=14, color=WHITE, bold=True,
             font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, out_x, out_y + 580_000, out_w, 380_000,
             ".rknn / .rkllm\nSHA-256 + RSA sign (optional)", size=9.5, color=WHITE,
             italic=True, align=PP_ALIGN.CENTER, line_spacing=1.35)

    # Bottom caption
    add_text(s, 660_000, SLIDE_H.emu - 820_000, SLIDE_W.emu - 1_320_000, 280_000,
             "缺 SDK 时 cluster-level 精细降级：always-available 永远可用；"
             "rknn / rkllm 仅失败的那簇跳过（logger.warning + skip），其他正常运行",
             size=10, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, repo="pet-quantize")


def slide_quantize_dualmode():
    s = new_slide()
    draw_tech_header(s, "pet-quantize",
                     "pet-quantize · Dual-mode Runner",
                     "同一类 · 两种初始化 · 被 pet-eval 跨仓引用")
    # Two side-by-side panels: PC-sim vs on-device
    panel_w = (SLIDE_W.emu - 1_640_000 - 300_000) // 2
    panel_h = 3_400_000
    panel_y = 1_800_000
    left_x = 820_000
    right_x = left_x + panel_w + 300_000

    panels = [
        {
            "title": "PC 模拟模式",
            "badge": "dev / CI · 无硬件",
            "color": BLUE,
            "init":  "RKLLMRunner(\n    model_path=uri,\n    target=None,\n    device_id=None,\n)",
            "note":  "PC 端 Python 模拟\n用于 orchestrator / pet-eval\n本地联调 · 不需要 rk3576",
        },
        {
            "title": "端侧 ADB 模式",
            "badge": "prod · on-device rk3576",
            "color": AMBER,
            "init":  "RKLLMRunner(\n    model_path=uri,\n    target=\"rk3576\",\n    device_id=\"ADB_SERIAL\",\n)",
            "note":  "通过 ADB 桥接真机\n代码不变 · flip 参数即可切换\nPhase 5 硬件接入触发",
        },
    ]
    for i, P in enumerate(panels):
        x = left_x if i == 0 else right_x
        add_rect(s, x, panel_y, panel_w, panel_h, fill=WHITE, line=P["color"], line_w=1.2)
        add_rect(s, x, panel_y, panel_w, 100_000, fill=P["color"])
        add_text(s, x, panel_y + 200_000, panel_w, 500_000,
                 P["title"], size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, x, panel_y + 750_000, panel_w, 300_000,
                 P["badge"], size=11, color=P["color"], bold=True, italic=True, align=PP_ALIGN.CENTER)
        # code init
        code_y = panel_y + 1_200_000
        add_rect(s, x + 280_000, code_y, panel_w - 560_000, 1_400_000, fill=GRAY_100)
        add_text(s, x + 380_000, code_y + 120_000, panel_w - 760_000, 1_200_000,
                 P["init"], size=11, color=NAVY, font=FONT_MONO, line_spacing=1.15)
        # note
        add_text(s, x + 280_000, panel_y + panel_h - 900_000, panel_w - 560_000, 800_000,
                 P["note"], size=11, color=GRAY_700, italic=True, line_spacing=1.5, align=PP_ALIGN.CENTER)

    # Bottom: cross-repo consumer
    add_text(s, 660_000, SLIDE_H.emu - 820_000, SLIDE_W.emu - 1_320_000, 280_000,
             "pet-eval QuantizedVlmEvaluator 通过 lazy `from pet_quantize.inference.rkllm_runner import RKLLMRunner` 跨仓引用；"
             "硬件到位前即可本地跑通全路径",
             size=10, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s, repo="pet-quantize")


# ---- pet-ota (2 slides) ----

def slide_ota_tech():
    s = new_slide()
    draw_tech_header(s, "pet-ota",
                     "pet-ota · 发布 + Canary Rollout 技术栈",
                     "bsdiff4 差分 · boto3/requests 多 backend · durable state · 可选 RSA 签名")
    techs = [
        ("bsdiff4",            "二进制 delta 差分更新\n@retry(3) tenacity 大文件 IO flake", BLUE),
        ("boto3 + STORAGE",    "S3 backend 发布\nURI scheme dispatch (file/s3/http)",       TEAL),
        ("requests PUT",       "HTTP backend · bearer/basic/no-auth\n自托管服务器",          AMBER),
        ("SHA-256 manifest",   "每 tarball 校验\nupload 前必过",                             REPO_COLORS["pet-ota"]),
        ("RSA signing (lazy)", "from pet_quantize.packaging.verify_package\nsoft-fail on ImportError", CORAL),
        ("Durable FSM state",  "deployments/<id>.json 续跑\n48h observation 崩溃可恢复",     GRAY_700),
    ]
    draw_tech_stack_grid(s, 1_800_000, techs, cols=3, card_h=1_550_000, gap=220_000)
    add_footer(s, repo="pet-ota")


def slide_ota_fsm():
    s = new_slide()
    draw_tech_header(s, "pet-ota",
                     "pet-ota · Canary Rollout 5-state FSM",
                     "Happy path 5 step · 任何 step 失败 → ROLLING_BACK → ROLLED_BACK · resume-from-state")
    # Happy path states horizontal
    happy = [
        ("GATE_CHECK",         "5 checks from params\nmin_dpo / min_days / ...",  BLUE),
        ("CANARY_DEPLOYING",   "canary_percentage=5%\n首批部署 canary group",        TEAL),
        ("CANARY_OBSERVING",   "observe=48h\nfailure_rate 监控",                   AMBER),
        ("FULL_DEPLOYING",     "全量 rollout\n剩余 95% devices",                    REPO_COLORS["pet-ota"]),
        ("DONE",               "rollout 完成\nmodel in production",                 TEAL),
    ]
    state_y = 1_900_000
    state_w = 1_880_000
    state_h = 1_100_000
    gap = 140_000
    total_w = 5 * state_w + 4 * gap
    start_x = (SLIDE_W.emu - total_w) // 2
    state_positions = []
    for i, (label, sub, color) in enumerate(happy):
        x = start_x + i * (state_w + gap)
        state_positions.append((x, state_y, state_w, state_h))
        add_rect(s, x, state_y, state_w, state_h, fill=WHITE, line=color, line_w=1.5)
        add_rect(s, x, state_y, state_w, 70_000, fill=color)
        add_text(s, x, state_y + 160_000, state_w, 420_000,
                 label, size=11.5, color=NAVY, bold=True, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)
        add_text(s, x + 80_000, state_y + 600_000, state_w - 160_000, state_h - 660_000,
                 sub, size=9.5, color=GRAY_500, align=PP_ALIGN.CENTER, line_spacing=1.4)
        if i < len(happy) - 1:
            add_arrow_right(s, x + state_w + 20_000, state_y + state_h // 2 - 80_000,
                            gap - 40_000, 160_000, color=BLUE)

    # Rollback row (below states 2/3/4)
    rb_y = state_y + state_h + 1_000_000
    rb_states = [
        ("ROLLING_BACK",  "backend.abort()\nrollback_timeout=5min", CORAL),
        ("ROLLED_BACK",   "status updated\nalert emitted",           GRAY_500),
    ]
    # Position under states 3 and 4
    rb_x_start = state_positions[2][0] + state_w + gap
    rb_w = state_w
    for i, (label, sub, color) in enumerate(rb_states):
        x = rb_x_start + i * (rb_w + gap)
        add_rect(s, x, rb_y, rb_w, state_h, fill=WHITE, line=color, line_w=1.5)
        add_rect(s, x, rb_y, rb_w, 70_000, fill=color)
        add_text(s, x, rb_y + 160_000, rb_w, 420_000,
                 label, size=11.5, color=NAVY, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, x + 80_000, rb_y + 600_000, rb_w - 160_000, state_h - 660_000,
                 sub, size=9.5, color=GRAY_500, align=PP_ALIGN.CENTER, line_spacing=1.4)

    # Draw downward rollback-trigger arrows from states 2/3/4 to ROLLING_BACK
    rb_target_x = rb_x_start + rb_w // 2
    for trigger_idx in (1, 2, 3):
        src_x, src_y, w, h = state_positions[trigger_idx]
        # Short vertical drop with coral color
        drop_x = src_x + w // 2 - 5_000
        drop_start_y = src_y + h
        drop_end_y = rb_y
        # draw as a thin vertical bar + small arrow
        add_rect(s, drop_x, drop_start_y + 40_000, 10_000,
                 drop_end_y - drop_start_y - 80_000, fill=CORAL)
        # arrowhead triangle
        arr = new_slide_arrow_head(s, drop_x - 30_000, drop_end_y - 40_000,
                                   70_000, 60_000, color=CORAL, direction="down")

    # Resume annotation
    add_rect(s, 660_000, SLIDE_H.emu - 820_000, SLIDE_W.emu - 1_320_000, 380_000,
             fill=GRAY_100)
    add_text(s, 820_000, SLIDE_H.emu - 760_000, SLIDE_W.emu - 1_640_000, 300_000,
             "Resume-from-state · 任一中断状态（canary_deploying / canary_observing / full_deploying）"
             "都会从 deployments/<id>.json 续跑 · crash / restart 不丢进度",
             size=10, color=GRAY_700, italic=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, repo="pet-ota")


def new_slide_arrow_head(s, x, y, w, h, *, color=BLUE, direction="down"):
    """Tiny arrow head drawn as a rectangle tip marker."""
    shape = MSO_SHAPE.DOWN_ARROW if direction == "down" else MSO_SHAPE.RIGHT_ARROW
    arr = s.shapes.add_shape(shape, Emu(x), Emu(y), Emu(w), Emu(h))
    arr.fill.solid(); arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    arr.shadow.inherit = False
    return arr


# ---- pet-id (2 slides) ----

def slide_id_tech():
    s = new_slide()
    draw_tech_header(s, "pet-id",
                     "pet-id · 独立 CLI 工具 · 核心技术栈",
                     "spec §5.2 零 pet-* 运行时依赖 · 仅通过 matrix row 做版本对齐")
    techs = [
        ("ultralytics YOLOv10",  "pet detector\nCOCO 15=cat, 16=dog",                BLUE),
        ("torchreid OSNet x0_25", "ReID embedding\n512 dim + L2 normalize",           TEAL),
        ("sha256(L2-norm f32)",  "content-addressable pet_id\n跨 host 跨 dtype 确定", AMBER),
        ("Pydantic PetCard",     "独立数据模型\n不在 pet-schema",                     REPO_COLORS["pet-id"]),
        ("mmpose (optional)",    "[pose] extras\nAP-10K keypoints",                  CORAL),
        ("transformers Qwen2-VL",  "[narrative] extras\nbehavior 描述",               GRAY_700),
    ]
    draw_tech_stack_grid(s, 1_800_000, techs, cols=3, card_h=1_550_000, gap=220_000)
    add_footer(s, repo="pet-id")


def slide_id_flow():
    s = new_slide()
    draw_tech_header(s, "pet-id",
                     "pet-id · Register + Identify 双泳道",
                     "两条独立流程 · 共享 detector / embedder / 磁盘 library")
    # Two swimlanes: register (top) + identify (bottom)
    lane_h = 1_700_000
    lane_gap = 400_000
    reg_y = 1_900_000
    idn_y = reg_y + lane_h + lane_gap

    def draw_lane(y, color, label, badge, stages):
        # Label
        add_rect(s, 660_000, y, 1_200_000, lane_h, fill=color)
        add_text(s, 660_000, y + 280_000, 1_200_000, 400_000,
                 label, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, 660_000, y + 800_000, 1_200_000, 400_000,
                 badge, size=10, color=WHITE, italic=True, align=PP_ALIGN.CENTER, line_spacing=1.35)
        # Stages
        start_x = 660_000 + 1_200_000 + 200_000
        usable = SLIDE_W.emu - 660_000 - start_x
        n = len(stages)
        gap = 120_000
        box_w = (usable - gap * (n - 1)) // n
        for i, (lbl, sub) in enumerate(stages):
            x = start_x + i * (box_w + gap)
            add_rect(s, x, y, box_w, lane_h, fill=WHITE, line=color, line_w=1.0)
            add_rect(s, x, y, box_w, 60_000, fill=color)
            add_text(s, x, y + 180_000, box_w, 400_000,
                     lbl, size=11, color=NAVY, bold=True,
                     font=FONT_MONO, align=PP_ALIGN.CENTER)
            add_text(s, x + 60_000, y + 600_000, box_w - 120_000, lane_h - 660_000,
                     sub, size=9.5, color=GRAY_500, align=PP_ALIGN.CENTER, line_spacing=1.4)
            if i < n - 1:
                add_arrow_right(s, x + box_w + 20_000, y + lane_h // 2 - 80_000,
                                gap - 40_000, 160_000, color=GRAY_500)

    draw_lane(reg_y, BLUE, "Register", "petid register \n<photo / dir / video>", [
        ("YOLOv10",    "detect pet bbox"),
        ("crop",       "largest bbox area"),
        ("OSNet",      "embed → 512 dim"),
        ("L2 norm + sha256",   "pet_id[:8]"),
        ("library.save",  "PetCard +\ncrop + .npy"),
    ])
    draw_lane(idn_y, AMBER, "Identify", "petid identify \n<query.jpg>", [
        ("YOLOv10",         "detect all bbox"),
        ("crop each",       "per-detection"),
        ("OSNet embed",     "query vector"),
        ("library.identify", "cosine max\n≥ threshold"),
        ("→ result",        "{pet_id, name, score}"),
    ])
    add_footer(s, repo="pet-id")


# ---- Ecosystem summary (2 slides) ----

def slide_ecosystem_overview():
    s = new_slide()
    add_header_bar(s, "生态全景 · 技术联通图",
                   "pet-schema 契约 + pet-infra 运行时 + 7 流水线仓 + 1 独立工具 · 一图看全",
                   section_num="§ 生态闭环", accent=NAVY)
    # 9-repo layout: chain top, pet-infra band, pet-id floater
    chain = [
        ("pet-schema",     "Pydantic / Alembic",         REPO_COLORS["pet-schema"]),
        ("pet-data",       "FFmpeg / SQLite",            REPO_COLORS["pet-data"]),
        ("pet-annotation", "4 paradigm / Label Studio",  REPO_COLORS["pet-annotation"]),
        ("pet-train",      "LLaMA-Factory / PEFT",       REPO_COLORS["pet-train"]),
        ("pet-eval",       "HF / BERTScore / gate",      REPO_COLORS["pet-eval"]),
        ("pet-quantize",   "RKLLM / RKNN",               REPO_COLORS["pet-quantize"]),
        ("pet-ota",        "bsdiff4 / canary FSM",       REPO_COLORS["pet-ota"]),
    ]
    # Row 1 — chain
    row_y = 1_900_000
    box_w = 1_540_000
    box_h = 1_250_000
    gap = 120_000
    total_w = 7 * box_w + 6 * gap
    start_x = (SLIDE_W.emu - total_w) // 2
    for i, (name, tech, color) in enumerate(chain):
        x = start_x + i * (box_w + gap)
        add_rect(s, x, row_y, box_w, box_h, fill=WHITE, line=color, line_w=1.2)
        add_rect(s, x, row_y, box_w, 70_000, fill=color)
        add_text(s, x, row_y + 170_000, box_w, 360_000,
                 name, size=11, color=color, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, x + 60_000, row_y + 600_000, box_w - 120_000, box_h - 650_000,
                 tech, size=9, color=GRAY_700, align=PP_ALIGN.CENTER, line_spacing=1.4)
        # ModelCard arrow between
        if i < len(chain) - 1:
            add_arrow_right(s, x + box_w + 10_000, row_y + box_h // 2 - 70_000,
                            gap - 20_000, 140_000, color=GRAY_500)

    # pet-infra band (shared runtime)
    infra_y = row_y + box_h + 260_000
    infra_color = REPO_COLORS["pet-infra"]
    add_rect(s, start_x, infra_y, total_w, 550_000, fill=infra_color)
    add_text(s, start_x, infra_y + 90_000, total_w, 400_000,
             "pet-infra · shared runtime",
             size=13, color=WHITE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
    add_text(s, start_x, infra_y + 350_000, total_w, 250_000,
             "mmengine Registry · Hydra compose · networkx DAG · ClearML · entry_points 发现",
             size=10, color=GRAY_300, italic=True, align=PP_ALIGN.CENTER)
    # Vertical connectors from chain to infra band
    for i in range(len(chain)):
        x = start_x + i * (box_w + gap) + box_w // 2 - 5_000
        add_rect(s, x, row_y + box_h, 10_000, infra_y - row_y - box_h,
                 fill=infra_color)

    # pet-id independent chip
    pid_y = infra_y + 720_000
    pid_color = REPO_COLORS["pet-id"]
    pid_w = 3_600_000
    pid_x = (SLIDE_W.emu - pid_w) // 2
    add_rect(s, pid_x, pid_y, pid_w, 540_000, fill=WHITE, line=pid_color, line_w=1.5)
    add_rect(s, pid_x, pid_y, 80_000, 540_000, fill=pid_color)
    add_text(s, pid_x + 160_000, pid_y + 80_000, pid_w - 320_000, 300_000,
             "pet-id · 独立 CLI 工具", size=13, color=pid_color, bold=True, font=FONT_MONO)
    add_text(s, pid_x + 160_000, pid_y + 280_000, pid_w - 320_000, 240_000,
             "YOLOv10 · torchreid OSNet · sha256 content-address · 零 pet-* 运行时依赖",
             size=9.5, color=GRAY_700, italic=True)

    # Legend at bottom
    add_text(s, 660_000, SLIDE_H.emu - 600_000, SLIDE_W.emu - 1_320_000, 280_000,
             "箭头 = ModelCard 契约逐级传递 · 垂直连线 = pet-infra runtime 贯穿 7 仓 · "
             "下沿独立 chip = pet-id（不参与 pipeline 主链，仅 matrix 登记）",
             size=10, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
    add_footer(s)


def slide_ecosystem_principles():
    s = new_slide()
    add_header_bar(s, "设计哲学 · 为什么这套架构能闭环",
                   "5 条共通原则 · 覆盖依赖 / 发现 / 确定性 / 硬件 / 契约",
                   section_num="§ 生态原则", accent=NAVY)
    principles = [
        {
            "num": "01",
            "name": "契约单一真理源",
            "tech": "pet-schema (Pydantic + Alembic)",
            "body": "所有跨仓类型住在一个仓；\n8 仓 import pet_schema 即得；\n链首零上游 → 不会版本协商地狱",
            "color": BLUE,
        },
        {
            "num": "02",
            "name": "β peer-dep 依赖松耦合",
            "tech": "pet-infra runtime 贯穿",
            "body": "不放 pyproject.dependencies；\nCI step 1 显式装 matrix pin；\n一处 bump 不扩散 7 个下游 diff",
            "color": TEAL,
        },
        {
            "num": "03",
            "name": "Registry + @register_module",
            "tech": "7 registries + entry_points 发现",
            "body": "每仓 plugin 装即注册；\norchestrator registry.build(type);\n新增 plugin 无需改 pet-infra",
            "color": AMBER,
        },
        {
            "num": "04",
            "name": "内容寻址确定性",
            "tech": "sha256(config) · sha256(embedding)",
            "body": "stage_config_sha / card_id / pet_id；\n跨 host 跨 dtype 同输入 → 同输出；\nresume-from-cache + migration-safe",
            "color": REPO_COLORS["pet-quantize"],
        },
        {
            "num": "05",
            "name": "Dual-mode hardware",
            "tech": "PC sim ↔ on-device ADB",
            "body": "RKLLMRunner / RKNNRunner 单一 class；\ntarget=None 本地；target+device_id 真机；\n硬件未到位即可全路径联调",
            "color": CORAL,
        },
    ]
    # 5 cards in single row (tight)
    start_y = 1_900_000
    card_w = (SLIDE_W.emu - 1_320_000 - 4 * 200_000) // 5
    card_h = 3_800_000
    gap = 200_000
    start_x = 660_000
    for i, p in enumerate(principles):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, start_y, card_w, card_h, fill=WHITE, line=GRAY_300, line_w=0.75)
        add_rect(s, x, start_y, card_w, 100_000, fill=p["color"])
        # Big number
        add_text(s, x + 200_000, start_y + 220_000, card_w - 400_000, 500_000,
                 p["num"], size=30, color=p["color"], bold=True, font=FONT_MONO)
        # Name
        add_text(s, x + 200_000, start_y + 820_000, card_w - 400_000, 520_000,
                 p["name"], size=14, color=NAVY, bold=True, line_spacing=1.25)
        # Tech signature
        add_text(s, x + 200_000, start_y + 1_440_000, card_w - 400_000, 440_000,
                 p["tech"], size=9.5, color=p["color"], italic=True, bold=True, font=FONT_MONO, line_spacing=1.35)
        # Body
        add_text(s, x + 200_000, start_y + 2_000_000, card_w - 400_000, card_h - 2_100_000,
                 p["body"], size=10.5, color=GRAY_700, line_spacing=1.55)
    add_footer(s)


# ============================================================================
# Pitch-deck front matter (non-technical framing)
# ============================================================================

def slide_problem():
    """pitch-style: single bold statement + 3 concrete pain cards."""
    s = new_slide()
    # Full-bleed subtle background
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=GRAY_100)
    # Eyebrow
    add_text(s, 820_000, 700_000, 10_000_000, 400_000,
             "THE PROBLEM", size=12, color=BLUE, bold=True)
    # Hero statement
    add_text(s, 820_000, 1_150_000, 10_500_000, 1_500_000,
             "养宠家庭最大的智能化缺口：\n摄像头能看见，却看不懂。",
             size=34, color=NAVY, bold=True, line_spacing=1.25)

    # Three pain cards
    pains = [
        {
            "num": "01",
            "title": "监控盲区",
            "body": "传统 IP cam 只记录画面；是不是在吃饭 / 有没有呕吐 /\n情绪是否低落，需要主人自己看回放判断",
            "color": CORAL,
        },
        {
            "num": "02",
            "title": "异常事后发现",
            "body": "呕吐、拒食、行为异常往往延时几小时才被察觉；\n急症（异物 / 中毒 / 胰腺炎）黄金窗口错失",
            "color": AMBER,
        },
        {
            "num": "03",
            "title": "投喂与状态脱节",
            "body": "现有喂食器靠定时；宠物状态（刚吐完 / 正在吃 / 外出）\n无法反馈投喂策略。IoT ≠ 智能",
            "color": BLUE,
        },
    ]
    card_y = 3_350_000
    card_h = 2_700_000
    card_w = 3_500_000
    gap = 300_000
    start_x = (SLIDE_W.emu - card_w * 3 - gap * 2) // 2
    for i, p in enumerate(pains):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, card_y, card_w, card_h, fill=WHITE, line=GRAY_300, line_w=0.75)
        # Top color bar
        add_rect(s, x, card_y, card_w, 80_000, fill=p["color"])
        # Big number
        add_text(s, x + 280_000, card_y + 220_000, 1_500_000, 600_000,
                 p["num"], size=36, color=p["color"], bold=True, font=FONT_MONO)
        # Title
        add_text(s, x + 280_000, card_y + 860_000, card_w - 560_000, 450_000,
                 p["title"], size=18, color=NAVY, bold=True)
        # Body
        add_text(s, x + 280_000, card_y + 1_360_000, card_w - 560_000, card_h - 1_440_000,
                 p["body"], size=12, color=GRAY_700, line_spacing=1.5)
    add_footer(s)


def slide_market_why_now():
    """pitch-style: market size 国内+国外 + why now."""
    s = new_slide()
    add_header_bar(s, "市场概况 · 为什么是现在",
                   "宠物经济规模 × 端侧 AI 成熟期 = 新一代智能硬件机会窗",
                   section_num="MARKET & TIMING")

    # Two market columns
    start_y = 1_800_000
    col_h = 2_600_000
    col_w = 5_500_000
    gap = 300_000
    start_x = 660_000
    markets = [
        {
            "region": "中国",
            "color": CORAL,
            "headline": "千亿级别 · 宠物智能硬件年增双位数",
            "facts": [
                "养宠家庭渗透率 持续上升",
                "宠物经济年规模：千亿元量级",
                "智能宠物家居品类：年增 20 – 30%",
                "一线城市高端用户愿付费：健康 / 行为管理",
            ],
        },
        {
            "region": "海外",
            "color": TEAL,
            "headline": "北美/欧洲成熟市场 · AI 是下一增长轴",
            "facts": [
                "北美养宠渗透率 > 60%（全球最高）",
                "人均宠物支出数倍于亚洲",
                "Pet tech 板块获 VC 持续加注",
                "智能摄像头 / 猫砂盆已 commodity，AI 差异化是新维度",
            ],
        },
    ]
    for i, m in enumerate(markets):
        x = start_x + i * (col_w + gap)
        add_rect(s, x, start_y, col_w, col_h, fill=WHITE, line=GRAY_300)
        # Color strip left
        add_rect(s, x, start_y, 110_000, col_h, fill=m["color"])
        # Region title
        add_text(s, x + 280_000, start_y + 200_000, col_w - 400_000, 450_000,
                 m["region"], size=22, color=m["color"], bold=True)
        add_text(s, x + 280_000, start_y + 720_000, col_w - 400_000, 400_000,
                 m["headline"], size=13, color=GRAY_700, italic=True)
        # Facts
        add_bullets(s, x + 280_000, start_y + 1_280_000, col_w - 400_000, col_h - 1_350_000,
                    m["facts"], size=13, bullet_color=m["color"], color=NAVY)

    # Why now band — taller to hold 2 body lines cleanly
    why_y = start_y + col_h + 240_000
    add_rect(s, 660_000, why_y, SLIDE_W.emu - 1_320_000, 1_500_000, fill=NAVY)
    add_text(s, 820_000, why_y + 140_000, 6_000_000, 400_000,
             "WHY NOW · 三个因素第一次合流",
             size=12, color=BLUE, bold=True)
    triggers = [
        ("1. VLM 可用", "Qwen2-VL / LLaVA 等开源 VLM 2023-24 成熟\n可在端侧量化后跑行为理解"),
        ("2. 端侧芯片", "RK3576 / 海思 / 高通 等 NPU 量产\n4-6 TOPS 满足 VLM + 音频多模态"),
        ("3. 消费需求", "主人要离线 / 低延迟 / 隐私\n云端方案开始被拒绝"),
    ]
    trig_w = (SLIDE_W.emu - 1_640_000) // 3
    for i, (title, body) in enumerate(triggers):
        tx = 820_000 + i * trig_w
        add_text(s, tx, why_y + 520_000, trig_w - 120_000, 300_000,
                 title, size=13, color=BLUE, bold=True, font=FONT_MONO)
        add_text(s, tx, why_y + 830_000, trig_w - 120_000, 620_000,
                 body, size=10, color=GRAY_300, line_spacing=1.45)

    add_footer(s)


def slide_competitive_map():
    """2×2 positioning map: AI depth × deployment location."""
    s = new_slide()
    add_header_bar(s, "竞品定位 · 2×2 Positioning",
                   "X 轴 · AI 能力深度　|　Y 轴 · 部署位置",
                   section_num="COMPETITIVE LANDSCAPE")

    # Chart area
    chart_x = 2_400_000
    chart_y = 1_850_000
    chart_w = 7_400_000
    chart_h = 4_000_000

    # Chart bg + border
    add_rect(s, chart_x, chart_y, chart_w, chart_h, fill=WHITE, line=GRAY_500, line_w=1.0)
    # Midlines (dashed feel: thin gray)
    mid_x = chart_x + chart_w // 2
    mid_y = chart_y + chart_h // 2
    add_rect(s, mid_x, chart_y, 3_000, chart_h, fill=GRAY_300)
    add_rect(s, chart_x, mid_y, chart_w, 3_000, fill=GRAY_300)

    # Axis labels — use text prefix instead of emoji for font-safety
    # X-axis
    add_text(s, chart_x, chart_y + chart_h + 140_000, chart_w, 320_000,
             "AI 能力深度  →  基础检测 · · · RFID · · · 云端分类 · · · 端侧 VLM + 多模态",
             size=11, color=GRAY_500, italic=True, align=PP_ALIGN.CENTER)
    # Y-axis labels, ASCII-safe
    add_text(s, chart_x - 1_900_000, chart_y + 80_000, 1_800_000, 280_000,
             "[ 云端 ]  部署", size=11, color=GRAY_500, italic=True, align=PP_ALIGN.RIGHT, bold=True)
    add_text(s, chart_x - 1_900_000, chart_y + chart_h - 400_000, 1_800_000, 280_000,
             "[ 端侧 ]  部署", size=11, color=GRAY_500, italic=True, align=PP_ALIGN.RIGHT, bold=True)

    # Quadrant labels (subtle, in corners)
    add_text(s, chart_x + 80_000, chart_y + 80_000, 3_000_000, 300_000,
             "Q1  基础 IoT / 云 · · 红海", size=9, color=GRAY_500, bold=True)
    add_text(s, mid_x + 80_000, chart_y + 80_000, 3_000_000, 300_000,
             "Q2  云端 AI · · 隐私与延迟待解", size=9, color=GRAY_500, bold=True)
    add_text(s, chart_x + 80_000, mid_y + 80_000, 3_000_000, 300_000,
             "Q3  端侧但 AI 弱 · · 空档", size=9, color=GRAY_500, bold=True)
    add_text(s, mid_x + 80_000, mid_y + 80_000, 3_000_000, 300_000,
             "Q4  端侧 VLM · · 我们", size=10, color=BLUE, bold=True)

    # Plot competitors (x, y are 0-1 within chart)
    # Dots dispersed to avoid label collisions in dense quadrants.
    competitors = [
        # (label, x_frac, y_frac, color, size)
        # Q1 (云端 · 基础) — 4 foreign cam/feeder brands, staircase layout
        ("Furbo",        0.12, 0.18, GRAY_500, 10),
        ("Petcube",      0.28, 0.12, GRAY_500, 10),
        ("WOpet",        0.15, 0.30, GRAY_500, 10),
        ("Petnet",       0.32, 0.26, GRAY_500, 10),
        # Q2 (云端 · 较深 AI) — nobody today sits deep; placed at frontier
        # (omitted; cleaner visual)
        # Q3 (端侧 · 基础) — RFID + China smart-home brands, staircase
        ("Sure Petcare", 0.10, 0.58, GRAY_500, 10),
        ("凡米",          0.14, 0.72, GRAY_500, 10),
        ("米家",          0.22, 0.62, GRAY_500, 10),
        ("CATLINK",      0.28, 0.80, GRAY_500, 10),
        ("PETKIT 小佩",    0.36, 0.68, GRAY_500, 10),
        # Q4 (端侧 · VLM 级深 AI) — our position, isolated and emphasised
        ("Train-Pet-Pipeline", 0.80, 0.78, BLUE, 13),
    ]
    for label, xf, yf, color, fsize in competitors:
        is_us = label == "Train-Pet-Pipeline"
        # Dot
        dot_d = 240_000 if is_us else 160_000
        px = chart_x + int(xf * chart_w) - dot_d // 2
        py = chart_y + int(yf * chart_h) - dot_d // 2
        add_rect(s, px, py, dot_d, dot_d,
                 fill=color if not is_us else BLUE,
                 shape=MSO_SHAPE.OVAL)
        if is_us:
            # Outer ring for emphasis
            ring = 360_000
            add_rect(s, chart_x + int(xf * chart_w) - ring // 2,
                     chart_y + int(yf * chart_h) - ring // 2,
                     ring, ring,
                     fill=WHITE, line=BLUE, line_w=2, shape=MSO_SHAPE.OVAL)
            # Redraw dot on top
            add_rect(s, px, py, dot_d, dot_d, fill=BLUE, shape=MSO_SHAPE.OVAL)
        # Label
        lbl_x = chart_x + int(xf * chart_w) + dot_d // 2 + 50_000
        lbl_y = chart_y + int(yf * chart_h) - 160_000
        add_text(s, lbl_x, lbl_y, 2_400_000, 260_000,
                 label,
                 size=fsize if not is_us else 11,
                 color=BLUE if is_us else GRAY_700,
                 bold=is_us,
                 font=FONT_SANS if is_us else FONT_SANS)

    # (Caveat removed — subtitle already sets the qualitative-positioning framing.)
    add_footer(s)


def slide_competitor_foreign():
    """Table: foreign main competitors."""
    s = new_slide()
    add_header_bar(s, "国外主流竞品 · 能力矩阵",
                   "5 个代表性玩家 × 6 维度 · 数据来自公开产品页 + 用户评测",
                   section_num="FOREIGN COMPETITION")

    headers = ["产品", "定位", "AI 能力", "端侧 / 云", "多模态", "价格带"]
    rows = [
        ["Furbo",         "狗主摄像头 + 零食投掷",   "云端 bark 检测 / 宠物识别",   "云端为主",      "视频 + 零食",        "¥1500 – 3000"],
        ["Petcube",       "互动摄像头 + 激光逗猫",   "云端基础识别",                "云端",          "视频 + 双向音频",     "¥1000 – 2500"],
        ["Petnet",        "智能喂食器 (初创)",        "基本 IoT + 云端调度",         "云端",          "仅投喂",             "¥800 – 2000"],
        ["Sure Petcare",  "Microchip 识别喂食",       "RFID 硬件识别",               "端侧 · RFID",    "仅 RFID + 秤",       "¥1200 – 3000"],
        ["WOpet",         "喂食器 + 摄像头捆绑",      "基础动物检测",                "云端",          "视频 + 投喂",         "¥600 – 1500"],
    ]
    _draw_comp_table(s, headers, rows, y_start=1_820_000)

    # Takeaway band
    tk_y = SLIDE_H.emu - 1_000_000
    add_rect(s, 660_000, tk_y, SLIDE_W.emu - 1_320_000, 480_000, fill=GRAY_100)
    add_text(s, 820_000, tk_y + 120_000, SLIDE_W.emu - 1_640_000, 280_000,
             "共性 · AI 层级普遍停留在云端基础检测；端侧 VLM 级别行为理解 = 空白市场",
             size=12, color=NAVY, bold=True)
    add_footer(s)


def slide_competitor_china():
    """Table: China competitors."""
    s = new_slide()
    add_header_bar(s, "国内主流竞品 · 能力矩阵",
                   "5 个代表性玩家 × 6 维度 · 国产智能宠物硬件头部阵营",
                   section_num="CHINA COMPETITION")

    headers = ["产品", "定位", "AI 能力", "端侧 / 云", "多模态", "价格带"]
    rows = [
        ["PETKIT 小佩",   "全品类智能宠物硬件",       "摄像头基础识别 + 云端",       "云端为主",      "视频 + 投喂 + 饮水",    "¥500 – 2500"],
        ["米家 / 小米",   "米家生态 IoT 产品",        "基础动作识别",                "云端",          "视频 + 投喂",           "¥200 – 1000"],
        ["CATLINK 凯特灵", "智能猫砂 + 投喂",         "基础检测 / 数据统计",         "云端",          "猫砂 + 体重 + 投喂",    "¥800 – 2500"],
        ["凡米 Fami",     "喂食器专精品牌",           "定时 + 基础识别",             "云端",          "仅投喂",                "¥300 – 1200"],
        ["米家生态链 / 其他", "多个白牌",             "基础 IoT",                    "云端",          "视频 + 投喂",           "¥200 – 800"],
    ]
    _draw_comp_table(s, headers, rows, y_start=1_820_000)

    tk_y = SLIDE_H.emu - 1_000_000
    add_rect(s, 660_000, tk_y, SLIDE_W.emu - 1_320_000, 480_000, fill=GRAY_100)
    add_text(s, 820_000, tk_y + 120_000, SLIDE_W.emu - 1_640_000, 280_000,
             "共性 · 硬件创新领先，AI 算法栈薄弱；行为理解 / 音频事件 / 结构化报告是下一代竞争点",
             size=12, color=NAVY, bold=True)
    add_footer(s)


def _draw_comp_table(s, headers, rows, y_start):
    """Shared competitor-table renderer.

    col_widths total = 10_700_000 EMU = 11.70"; slide usable = 11.89"
    (13.33 - 2*0.722" margin). Leaves ~0.10" safety on each side."""
    col_widths = [1_500_000, 2_150_000, 2_300_000, 1_500_000, 1_650_000, 1_600_000]
    x0 = (SLIDE_W.emu - sum(col_widths)) // 2
    row_h = 500_000
    header_h = 500_000
    # Header
    cx = x0
    for i, (head, w) in enumerate(zip(headers, col_widths)):
        add_rect(s, cx, y_start, w, header_h, fill=NAVY)
        add_text(s, cx + 120_000, y_start, w - 240_000, header_h,
                 head, size=11, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        cx += w
    # Rows
    for ri, row in enumerate(rows):
        y = y_start + header_h + ri * row_h
        fill = GRAY_100 if ri % 2 else WHITE
        cx = x0
        for ci, (val, w) in enumerate(zip(row, col_widths)):
            add_rect(s, cx, y, w, row_h, fill=fill, line=GRAY_300, line_w=0.4)
            bold = ci == 0
            color = NAVY if ci == 0 else GRAY_700
            add_text(s, cx + 120_000, y, w - 240_000, row_h,
                     val, size=10.5, color=color, bold=bold,
                     anchor=MSO_ANCHOR.MIDDLE,
                     align=PP_ALIGN.LEFT if ci <= 1 else PP_ALIGN.LEFT)
            cx += w


def slide_differentiation():
    """Side-by-side: majority of competitors vs us."""
    s = new_slide()
    add_header_bar(s, "我们的差异化",
                   "唯一在端侧落地 VLM-grade 行为理解 + 音频事件 + 签名 OTA 的方案",
                   section_num="DIFFERENTIATION")

    # 3-col grid: [dim tag | 大多数竞品 | Train-Pet-Pipeline]
    # All within slide: slide 13.33" - 2×0.72" margin = 11.89" usable.
    # tag = 1.4"  + gap 0.25" + col 4.95" + gap 0.25" + col 4.95" ≈ 11.80"
    start_y = 1_700_000
    col_h = 4_500_000
    gap = 230_000
    tag_w = 1_280_000
    col_w = (SLIDE_W.emu - 1_320_000 - tag_w - 2 * gap) // 2
    tag_x = 660_000
    left_x = tag_x + tag_w + gap
    right_x = left_x + col_w + gap

    dimensions = [
        ("AI 能力",    "云端基础动作识别",               "端侧 VLM 行为理解 + 结构化 JSON"),
        ("多模态",     "视频为主，音频不做处理",         "视频 + 音频事件 (呕吐/吃饭/饮水/…)"),
        ("隐私",       "视频传云；隐私合规压力",         "全程端侧推理，视频不出设备"),
        ("延迟",       "2 – 5 秒起（往返云）",            "< 4 秒 P95（本地）"),
        ("可升级",     "厂商云服务耦合",                  "签名 OTA + canary 灰度 + 回滚"),
        ("离线可用",   "断网基本失效",                    "纯端侧能力；网络只用于 OTA"),
    ]

    # Headers
    # tag column header (empty placeholder)
    add_text(s, tag_x, start_y + 220_000, tag_w, 400_000,
             "维度", size=11, color=GRAY_500, italic=True, bold=True,
             align=PP_ALIGN.RIGHT)
    # left header — majority
    add_rect(s, left_x, start_y, col_w, 90_000, fill=GRAY_500)
    add_text(s, left_x + 180_000, start_y + 160_000, col_w - 360_000, 450_000,
             "大多数竞品", size=17, color=GRAY_500, bold=True)
    add_text(s, left_x + 180_000, start_y + 660_000, col_w - 360_000, 300_000,
             "基础 IoT + 云端", size=11, color=GRAY_500, italic=True)
    # right header — us
    add_rect(s, right_x, start_y, col_w, 90_000, fill=BLUE)
    add_text(s, right_x + 180_000, start_y + 160_000, col_w - 360_000, 450_000,
             "Train-Pet-Pipeline", size=17, color=BLUE, bold=True)
    add_text(s, right_x + 180_000, start_y + 660_000, col_w - 360_000, 300_000,
             "端侧 VLM + 多模态 + 签名 OTA", size=11, color=BLUE, italic=True)

    # Dimension rows
    rows_y = start_y + 1_100_000
    row_h = (col_h - 1_200_000) // len(dimensions)
    for ri, (dim, them, us) in enumerate(dimensions):
        y = rows_y + ri * row_h
        # Dimension tag — inside slide, right-aligned
        add_text(s, tag_x, y + 60_000, tag_w - 80_000, row_h - 80_000,
                 dim, size=12, color=NAVY, italic=False, bold=True,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Left cell
        add_rect(s, left_x, y + 40_000, col_w, row_h - 80_000,
                 fill=WHITE, line=GRAY_300, line_w=0.5)
        add_text(s, left_x + 180_000, y + 60_000, col_w - 360_000, row_h - 80_000,
                 them, size=11, color=GRAY_700, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
        # Right cell (accent left border)
        add_rect(s, right_x, y + 40_000, col_w, row_h - 80_000,
                 fill=WHITE, line=BLUE, line_w=1)
        add_rect(s, right_x, y + 40_000, 60_000, row_h - 80_000, fill=BLUE)
        add_text(s, right_x + 180_000, y + 60_000, col_w - 360_000, row_h - 80_000,
                 us, size=11, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)

    add_footer(s)


def slide_solution():
    """Big-picture solution: 3-pillar + hardware diagram + stats strip."""
    s = new_slide()
    add_header_bar(s, "解决方案",
                   "一台端侧 AI 设备 · 三大能力柱 · 签名 OTA 闭环",
                   section_num="OUR SOLUTION")

    # Three pillar cards (reused structure from earlier stats-card pattern)
    base_y = 1_800_000
    card_h = 1_700_000
    card_w = 3_450_000
    gap = 320_000
    start_x = (SLIDE_W.emu - card_w * 3 - gap * 2) // 2
    pillars = [
        {
            "icon": "👁",
            "title": "VLM 视觉语言",
            "body": "Qwen2-VL-2B LoRA SFT / DPO\n结构化 JSON (schema_compliance 0.99)\n端侧 RKLLM W4A16",
            "color": BLUE,
        },
        {
            "icon": "🎧",
            "title": "音频事件",
            "body": "PANNs MobileNetV2 AudioSet CNN\n5 类（呕吐 / 吃饭 / 饮水 / 环境 / 其他）\n端侧 RKNN INT8",
            "color": TEAL,
        },
        {
            "icon": "📦",
            "title": "签名 OTA",
            "body": "bsdiff4 差分更新\n5 状态 canary rollout + 48h 观察窗\nSHA-256 manifest + optional RSA 签名",
            "color": AMBER,
        },
    ]
    for i, p in enumerate(pillars):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, base_y, card_w, card_h, fill=WHITE, line=GRAY_300, line_w=0.75)
        add_rect(s, x, base_y, card_w, 110_000, fill=p["color"])
        add_text(s, x + 260_000, base_y + 250_000, 800_000, 600_000,
                 p["icon"], size=28, color=p["color"])
        add_text(s, x + 1_100_000, base_y + 280_000, card_w - 1_300_000, 500_000,
                 p["title"], size=18, color=NAVY, bold=True)
        add_text(s, x + 260_000, base_y + 900_000, card_w - 500_000, card_h - 1_000_000,
                 p["body"], size=11.5, color=GRAY_700, line_spacing=1.5)

    # Hardware side-bar — taller to hold 2-line desc per component
    hw_y = base_y + card_h + 280_000
    add_rect(s, 660_000, hw_y, SLIDE_W.emu - 1_320_000, 1_500_000, fill=NAVY)
    add_text(s, 820_000, hw_y + 140_000, 5_000_000, 400_000,
             "HARDWARE · 设备内置", size=12, color=BLUE, bold=True)
    # Use ASCII-safe glyphs so LibreOffice / non-emoji fonts render consistently
    components = [
        ("▣", "RK3576 SoC", "6 TOPS NPU · VLM + Audio\n同时推理"),
        ("▦", "Camera",     "1080p · 广角\n端侧流入 VLM"),
        ("◉", "Microphone", "高采样率 mic\n音频事件检测"),
        ("◈", "Feeder Motor", "定量投喂\n受 VLM 决策驱动"),
    ]
    cw = (SLIDE_W.emu - 1_640_000) // 4
    for i, (icon, name, desc) in enumerate(components):
        cx = 820_000 + i * cw
        add_text(s, cx, hw_y + 550_000, 400_000, 400_000,
                 icon, size=22, color=BLUE, bold=True)
        add_text(s, cx + 440_000, hw_y + 560_000, cw - 500_000, 300_000,
                 name, size=12, color=WHITE, bold=True, font=FONT_MONO)
        add_text(s, cx + 440_000, hw_y + 840_000, cw - 500_000, 600_000,
                 desc, size=9.5, color=GRAY_300, line_spacing=1.45)

    add_footer(s)


def slide_user_journey():
    """5-step user flow diagram."""
    s = new_slide()
    add_header_bar(s, "用户旅程",
                   "从开箱到日常预警 · 5 步体验",
                   section_num="USER JOURNEY")

    steps = [
        {
            "num": "1",
            "title": "开箱",
            "body": "扫码联 Wi-Fi\n3 分钟完成绑定",
            "color": BLUE,
        },
        {
            "num": "2",
            "title": "注册宠物",
            "body": "录 5-10 秒视频绕宠物一圈\n→ PetCard 生成 (pet-id)",
            "color": TEAL,
        },
        {
            "num": "3",
            "title": "日常监控",
            "body": "VLM 每帧读画面\n音频 CNN 听环境声\n结构化 JSON 上传 App",
            "color": AMBER,
        },
        {
            "num": "4",
            "title": "异常预警",
            "body": "呕吐 / 拒食 / 行为异常\napp 即时通知 + 回看视频\n前 30 秒预览免费",
            "color": CORAL,
        },
        {
            "num": "5",
            "title": "固件升级",
            "body": "签名 OTA canary 推送\n5% → 48h 观察 → 全量\n失败自动回滚",
            "color": REPO_COLORS["pet-ota"],
        },
    ]
    # Horizontal flow with arrows
    box_w = 2_000_000
    box_h = 2_600_000
    gap = 260_000
    total_w = len(steps) * box_w + (len(steps) - 1) * gap
    start_x = (SLIDE_W.emu - total_w) // 2
    row_y = 2_100_000

    for i, st in enumerate(steps):
        x = start_x + i * (box_w + gap)
        # Card
        add_rect(s, x, row_y, box_w, box_h, fill=WHITE, line=GRAY_300, line_w=0.75)
        # Color top band
        add_rect(s, x, row_y, box_w, 100_000, fill=st["color"])
        # Step circle
        circle_d = 560_000
        circle_x = x + box_w // 2 - circle_d // 2
        circle_y = row_y + 200_000
        add_rect(s, circle_x, circle_y, circle_d, circle_d,
                 fill=st["color"], shape=MSO_SHAPE.OVAL)
        add_text(s, circle_x, circle_y, circle_d, circle_d,
                 st["num"], size=26, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_MONO)
        # Title
        add_text(s, x + 100_000, row_y + 900_000, box_w - 200_000, 400_000,
                 st["title"], size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        # Body
        add_text(s, x + 200_000, row_y + 1_400_000, box_w - 400_000, box_h - 1_500_000,
                 st["body"], size=10.5, color=GRAY_700, align=PP_ALIGN.CENTER, line_spacing=1.5)
        # Arrow between
        if i < len(steps) - 1:
            arrow_x = x + box_w + 20_000
            arrow_y = row_y + box_h // 2 - 100_000
            add_arrow_right(s, arrow_x, arrow_y, gap - 40_000, 200_000, color=GRAY_300)

    add_footer(s)


# ============================================================================
# Content composition
# ============================================================================

# Front-matter: pitch-deck style
slide_title()
slide_problem()
slide_market_why_now()
slide_competitive_map()
slide_competitor_foreign()
slide_competitor_china()
slide_differentiation()
slide_solution()
slide_user_journey()

# Transition to technical
slide_pipeline_flow()
slide_version_matrix()

# ---- Section A: Contract + Runtime ----

slide_section_divider("A", "契约与运行时", "pet-schema · pet-infra")
slide_schema_tech()
slide_schema_flow()
slide_infra_tech()
slide_infra_flow()
slide_infra_registries()

# ---- Section B: Data + Annotation ----

slide_section_divider("B", "数据流", "pet-data · pet-annotation")
slide_data_tech()
slide_data_flow()
slide_annotation_tech()
slide_annotation_flow()

# ---- Section C: Train + Eval ----

slide_section_divider("C", "训练与评估", "pet-train · pet-eval")
slide_train_tech()
slide_train_flow()
slide_eval_tech()
slide_eval_flow()

# ---- Section D: Edge + Delivery ----

slide_section_divider("D", "端侧与发布", "pet-quantize · pet-ota")
slide_quantize_tech()
slide_quantize_flow()
slide_quantize_dualmode()
slide_ota_tech()
slide_ota_fsm()

# ---- Section E: Independent Tool ----

slide_section_divider("E", "独立工具", "pet-id")
slide_id_tech()
slide_id_flow()

# ---- Ecosystem summary — closes the technical narrative ----

slide_ecosystem_overview()
slide_ecosystem_principles()

# ---- Section F: Cross-cutting ----

slide_section_divider("F", "横切关注点", "依赖治理 · CI · 北极星", "Slides 35 – 37")

# Dependency governance
def slide_dep_governance():
    s = new_slide()
    add_header_bar(s, "依赖治理：3 种 pin 模式",
                   "生态优化收敛：从 6 种 ad-hoc 统一到 3 种 disciplined",
                   section_num="§ 10 · 横切 / 依赖")

    # Three columns
    start_y = 1_800_000
    col_h = 3_600_000
    col_w = 3_500_000
    gap = 300_000
    start_x = (SLIDE_W.emu - (col_w * 3 + gap * 2)) // 2

    styles = [
        {
            "name": "α 硬 pin",
            "subtitle": "pyproject.dependencies 写 git URL @vX.Y.Z",
            "color": CORAL,
            "applies": "pet-eval / pet-quantize / pet-ota\n(signing extras excepted)",
            "reasoning": [
                "适合：稳定契约型依赖",
                "如 pet-schema (consumers 当作数据类型)",
                "每 tag bump 需显式改 pyproject",
                "优势：pip 直装即可，CI 简单",
            ],
        },
        {
            "name": "β peer-dep",
            "subtitle": "不在 pyproject.dependencies；CI Step-1 explicit install",
            "color": TEAL,
            "applies": "pet-infra / pet-data / pet-annotation\npet-train / pet-eval / pet-quantize / pet-ota\n(相对 pet-infra 都是 β)",
            "reasoning": [
                "适合：频繁迭代型 peer",
                "如 pet-infra (版本常改)",
                "delayed RuntimeError guard in register_all()",
                "优势：裸 import 轻量 (IDE/静态分析)",
            ],
        },
        {
            "name": "跨仓 plugin no-pin",
            "subtitle": "pyproject 写 `pet-xxx` 不带版本",
            "color": BLUE,
            "applies": "pet-eval → pet-train, pet-quantize\npet-ota[signing] → pet-quantize",
            "reasoning": [
                "适合：跨仓 runtime plugin dep",
                "版本由 matrix row 锁定",
                "CI 显式装矩阵版本",
                "优势：plugin 发现松耦合",
            ],
        },
    ]
    for i, st in enumerate(styles):
        x = start_x + i * (col_w + gap)
        add_rect(s, x, start_y, col_w, col_h, fill=WHITE, line=GRAY_300)
        # Header
        add_rect(s, x, start_y, col_w, 90_000, fill=st["color"])
        add_text(s, x + 160_000, start_y + 180_000, col_w - 320_000, 480_000,
                 st["name"], size=20, color=st["color"], bold=True)
        add_text(s, x + 160_000, start_y + 680_000, col_w - 320_000, 480_000,
                 st["subtitle"], size=10, color=GRAY_500, italic=True, line_spacing=1.3)
        # Applies
        add_rect(s, x + 160_000, start_y + 1_280_000, col_w - 320_000, 10_000, fill=GRAY_300)
        add_text(s, x + 160_000, start_y + 1_340_000, col_w - 320_000, 300_000,
                 "适用仓库", size=10, color=GRAY_500, bold=True)
        add_text(s, x + 160_000, start_y + 1_640_000, col_w - 320_000, 800_000,
                 st["applies"], size=11, color=NAVY, line_spacing=1.4, font=FONT_MONO)
        # Reasoning
        add_bullets(s, x + 160_000, start_y + 2_550_000, col_w - 320_000, 980_000,
                    st["reasoning"], size=10.5, bullet_color=st["color"])

    add_footer(s, slide_num=35)


# CI guards matrix
def slide_ci_guards():
    s = new_slide()
    add_header_bar(s, "CI Guard 矩阵",
                   "生态优化把 5 类 workflow 一致化到 9 仓",
                   section_num="§ 10 · 横切 / CI")

    repos = ["pet-schema", "pet-infra", "pet-data", "pet-annotation",
             "pet-train", "pet-eval", "pet-quantize", "pet-ota", "pet-id"]
    workflows = [
        ("ci.yml",                  "lint + mypy + pytest"),
        ("peer-dep-smoke.yml",      "独立装序 smoke"),
        ("no-wandb-residue.yml",    "positive-list 扫 \\bwandb\\b"),
        ("schema_guard.yml",        "pet-schema dispatch 全链"),
        ("cross-repo-smoke.yml",    "matrix row 装序验证"),   # display-shortened; file is cross-repo-smoke-install.yml
    ]
    # Presence grid: cell value = True/False or special
    grid = {
        # repo, workflow idx -> "Y"/"N"/"."
        "pet-schema": ["Y", ".", ".", "Y", "."],  # schema_guard source
        "pet-infra": ["Y", "Y", "Y", ".", "Y"],
        "pet-data": ["Y", "Y", ".", ".", "."],  # 未加 no-wandb 因从未触过 W&B
        "pet-annotation": ["Y", "Y", ".", ".", "."],
        "pet-train": ["Y", "Y", "Y", ".", "."],
        "pet-eval": ["Y", "Y", "Y", ".", "."],
        "pet-quantize": ["Y", "Y", "Y", ".", "."],
        "pet-ota": ["Y", "Y", "Y", ".", "."],
        "pet-id": ["Y", ".", "Y", ".", "."],  # 独立工具无 peer-dep-smoke
    }

    # Layout
    left_label_w = 1_700_000
    cell_w = 1_720_000
    cell_h = 380_000
    header_y = 1_800_000
    grid_start_x = 660_000 + left_label_w

    # Column headers
    for i, (wf, desc) in enumerate(workflows):
        x = grid_start_x + i * cell_w
        add_rect(s, x, header_y, cell_w, cell_h + 120_000, fill=NAVY)
        add_text(s, x + 80_000, header_y + 50_000, cell_w - 160_000, 240_000,
                 wf, size=9.5, color=WHITE, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        add_text(s, x + 80_000, header_y + 280_000, cell_w - 160_000, 220_000,
                 desc, size=8, color=GRAY_300, align=PP_ALIGN.CENTER, italic=True)

    # Rows
    for ri, repo in enumerate(repos):
        y = header_y + cell_h + 120_000 + ri * cell_h
        # Repo label
        color = REPO_COLORS[repo]
        add_rect(s, 660_000, y, left_label_w, cell_h,
                 fill=WHITE, line=GRAY_300, line_w=0.5)
        add_rect(s, 660_000, y, 100_000, cell_h, fill=color)
        add_text(s, 790_000, y, left_label_w - 150_000, cell_h,
                 repo, size=10, color=NAVY, bold=True, font=FONT_MONO,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Cells
        for ci in range(len(workflows)):
            x = grid_start_x + ci * cell_w
            mark = grid[repo][ci]
            fill = WHITE
            line = GRAY_300
            if mark == "Y":
                fill = RGBColor(0xE7, 0xF4, 0xEE)  # light teal bg
                line = TEAL
            add_rect(s, x, y, cell_w, cell_h, fill=fill, line=line, line_w=0.5)
            if mark == "Y":
                add_text(s, x, y, cell_w, cell_h, "✓",
                         size=16, color=TEAL, bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                add_text(s, x, y, cell_w, cell_h, "—",
                         size=13, color=GRAY_300, bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 660_000, SLIDE_H.emu - 820_000, SLIDE_W.emu - 1_320_000, 280_000,
             "✓ = 已接入；— = 该仓无需要（如 pet-id 独立无 peer-dep / pet-schema 无下游 peer-dep smoke）",
             size=9.5, color=GRAY_500, italic=True)
    add_footer(s, slide_num=36)


def slide_northstar():
    s = new_slide()
    add_header_bar(s, "北极星四维度",
                   "生态优化 Phase 10 retrospective · 维持或提升，无退化",
                   section_num="§ 10 · 横切 / North Star")

    dims = [
        {
            "name": "Pluggability",
            "tagline": "插件化程度",
            "score": 5,
            "evidence": [
                "7 registries 覆盖 TRAINERS · EVALUATORS · CONVERTERS · METRICS · DATASETS · STORAGE · OTA",
                "每仓插件 entry-point 注册; orchestrator 发现加载",
            ],
            "color": BLUE,
        },
        {
            "name": "Flexibility",
            "tagline": "配置灵活性",
            "score": 5,
            "evidence": [
                "所有数值从 params.yaml 读（no-hardcode 铁律）",
                "ExperimentRecipe.variations ablation sweep",
                "Hydra defaults-list + override 组合",
            ],
            "color": TEAL,
        },
        {
            "name": "Extensibility",
            "tagline": "扩展容易程度",
            "score": 5,
            "evidence": [
                "每仓 architecture.md §5 显式说明加 plugin 步骤",
                "3 种 pin 模式覆盖所有 dependency 形态",
                "cross-repo-smoke-install 保证扩展不破 matrix",
            ],
            "color": AMBER,
        },
        {
            "name": "Comparability",
            "tagline": "实验可比性",
            "score": 5,
            "evidence": [
                "ClearML 单一 tracker; W&B 物理移除",
                "ModelCard 全链贯穿，metrics 可跨 run 对比",
                "Replay 机制：ModelCard → 同一 config 重跑",
            ],
            "color": CORAL,
        },
    ]

    # 2×2 grid
    start_y = 1_800_000
    card_w = 5_400_000
    card_h = 1_800_000
    gap_x = 280_000
    gap_y = 240_000
    start_x = 660_000

    for i, dim in enumerate(dims):
        r = i // 2
        c = i % 2
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        # Card
        add_rect(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_300)
        # Color band
        add_rect(s, x, y, 120_000, card_h, fill=dim["color"])
        # Name
        add_text(s, x + 280_000, y + 160_000, card_w - 320_000, 400_000,
                 dim["name"], size=18, color=NAVY, bold=True)
        add_text(s, x + 280_000, y + 560_000, card_w - 320_000, 280_000,
                 dim["tagline"], size=11, color=GRAY_500, italic=True)
        # Score dots (5 filled)
        dot_y = y + 170_000
        dot_x0 = x + card_w - 1_500_000
        for di in range(5):
            filled = di < dim["score"]
            add_rect(s, dot_x0 + di * 250_000, dot_y, 200_000, 200_000,
                     fill=dim["color"] if filled else GRAY_300,
                     shape=MSO_SHAPE.OVAL)
        # Evidence bullets
        add_bullets(s, x + 280_000, y + 900_000, card_w - 500_000, card_h - 1_000_000,
                    dim["evidence"], size=10.5, bullet_color=dim["color"],
                    color=GRAY_700)

    add_footer(s, slide_num=37)


# Status + future
def slide_status_closeout():
    s = new_slide()
    add_header_bar(s, "Phase 10 生态优化收官",
                   "10 Phase · 9 repo · 单日 sprint · DoD 5/5",
                   section_num="§ 11 · 现状")

    # Horizontal phase timeline
    phases = [
        # (num, repo_key, display_name, version)
        ("1",  "pet-schema",     "pet-schema",  "pre"),
        ("2",  "pet-infra",      "pet-infra",   "v2.6.0"),
        ("3",  "pet-data",       "pet-data",    "v1.3.0"),
        ("4",  "pet-annotation", "pet-annot",   "v2.1.1"),
        ("5",  "pet-train",      "pet-train",   "v2.0.2"),
        ("6",  "pet-eval",       "pet-eval",    "v2.3.0"),
        ("7",  "pet-quantize",   "pet-quant",   "v2.1.0"),
        ("8",  "pet-ota",        "pet-ota",     "v2.2.0"),
        ("9",  "pet-id",         "pet-id",      "v0.2.0"),
        ("10", None,             "closeout",    "matrix\n2026.10"),
    ]
    timeline_y = 2_300_000
    n = len(phases)
    box_w = 950_000
    box_h = 380_000
    gap = 100_000
    total_w = n * box_w + (n - 1) * gap
    start_x = (SLIDE_W.emu - total_w) // 2
    # Backbone line
    add_rect(s, start_x, timeline_y + box_h // 2 - 10_000, total_w, 20_000, fill=BLUE)
    for i, (num, repo_key, display, ver) in enumerate(phases):
        x = start_x + i * (box_w + gap)
        color = REPO_COLORS.get(repo_key, BLUE) if repo_key else BLUE
        # Phase dot
        add_rect(s, x + box_w // 2 - 80_000, timeline_y + box_h // 2 - 80_000,
                 160_000, 160_000, fill=color, shape=MSO_SHAPE.OVAL)
        # Phase number above
        add_text(s, x, timeline_y - 380_000, box_w, 280_000,
                 f"Phase {num}", size=10, color=GRAY_500, bold=True, align=PP_ALIGN.CENTER)
        # Repo name below (abbreviated where long names would wrap)
        add_text(s, x, timeline_y + box_h + 120_000, box_w, 280_000,
                 display, size=10.5, color=NAVY, bold=True, font=FONT_MONO, align=PP_ALIGN.CENTER)
        # Version lower
        add_text(s, x, timeline_y + box_h + 440_000, box_w, 420_000,
                 ver, size=9.5, color=color, align=PP_ALIGN.CENTER, line_spacing=1.2)

    # Stat strip at bottom
    stat_y = 4_700_000
    add_rect(s, 660_000, stat_y, SLIDE_W.emu - 1_320_000, 900_000, fill=GRAY_100)
    stats = [
        ("14+", "feature PR 合并"),
        ("9", "dev→main sync"),
        ("4", "workflow bug 修"),
        ("~533 LOC", "pet-eval 死代码删"),
        ("0", "open PR 收官"),
    ]
    stat_w = (SLIDE_W.emu - 1_320_000) / len(stats)
    for i, (num, label) in enumerate(stats):
        sx = 660_000 + int(i * stat_w)
        add_text(s, sx, stat_y + 140_000, int(stat_w), 400_000,
                 num, size=26, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, sx, stat_y + 580_000, int(stat_w), 250_000,
                 label, size=10, color=GRAY_500, align=PP_ALIGN.CENTER)

    add_footer(s, slide_num=38)


def slide_future():
    s = new_slide()
    add_header_bar(s, "Phase 5+ 硬件 · Follow-ups",
                   "retrospective §8 · 9 条（3 条 Phase 4 遗留 + 6 条本轮新增）",
                   section_num="§ 11 · 未来")

    # Two-column layout — usable width ~11.89"; col 5.65" each + 0.44" gap + 2×0.72" margin ≈ 12.18"
    left_x = 660_000
    col_w = 5_170_000
    right_x = left_x + col_w + 400_000
    start_y = 1_800_000
    col_h = 4_200_000

    # Left — 硬件阻塞
    add_rect(s, left_x, start_y, col_w, 90_000, fill=CORAL)
    add_text(s, left_x + 200_000, start_y + 160_000, col_w - 400_000, 400_000,
             "硬件阻塞项 · Phase 5 gate",
             size=16, color=CORAL, bold=True)
    hw_items = [
        ("RK3576 单元 + self-hosted runner", "Phase 5 核心前提；解锁所有 --hardware 路径"),
        ("quantize_validate.yml 真跑", "Phase 7 修了死调用；硬件到位第一步"),
        ("canary rollout 真机测试", "pet-ota FSM 本地 sim 已验；真机 48h 观察窗"),
        ("RKNN/RKLLM 实测 latency P95", "当前 matrix 有 threshold 无测量值"),
    ]
    add_bullets(s, left_x + 200_000, start_y + 640_000, col_w - 400_000, col_h - 700_000,
                hw_items, size=12, bullet_color=CORAL)

    # Right — 可并行推进
    add_rect(s, right_x, start_y, col_w, 90_000, fill=AMBER)
    add_text(s, right_x + 200_000, start_y + 160_000, col_w - 400_000, 400_000,
             "可并行推进 · 不等硬件",
             size=16, color=AMBER, bold=True)
    sw_items = [
        ("pet-schema __version__ 属性", "Phase 10 修 workflow 时发现缺失；加 + 同 parity test"),
        ("SFT→DPO composed recipe", "当前只有 SFT 顶层 recipe；拼两 stage 5 分钟事"),
        ("pet-ota signing-smoke CI", "跑一次带 [signing] extras 装矩阵版 pet-quantize"),
        ("pet-id heavy-backends CI", "pose + narrative + tracker extras 独立 job"),
        ("pet-eval _FALLBACK_OUTPUT dynamic", "从 pet_schema defaults 动态生成"),
    ]
    add_bullets(s, right_x + 200_000, start_y + 640_000, col_w - 400_000, col_h - 700_000,
                sw_items, size=12, bullet_color=AMBER)

    add_footer(s, slide_num=39)


# Closing slide
def slide_close():
    s = new_slide()
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    add_text(s, 820_000, 2_600_000, 11_000_000, 1_000_000,
             "一切就绪。", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 820_000, 3_700_000, 11_000_000, 600_000,
             "9 repos · tagged · main CI green · zero open PR",
             size=16, color=GRAY_300, align=PP_ALIGN.CENTER, italic=True)
    add_rect(s, SLIDE_W.emu // 2 - 800_000, 4_600_000, 1_600_000, 8_000, fill=BLUE)
    add_text(s, 820_000, 4_800_000, 11_000_000, 400_000,
             "pet-infra/docs/architecture/OVERVIEW.md · compatibility_matrix.yaml: 2026.10-ecosystem-cleanup",
             size=11, color=BLUE, align=PP_ALIGN.CENTER, font=FONT_MONO)


# ---- Compose everything ----
slide_dep_governance()
slide_ci_guards()
slide_northstar()
slide_status_closeout()
slide_future()
slide_close()


# Save
out = Path(__file__).parent / "Train-Pet-Pipeline-Technical-Overview.pptx"
prs.save(str(out))
print(f"Wrote {out}  ({out.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
